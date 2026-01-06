import os
import json
import requests
import re
import hashlib
import sqlite3
import time
from datetime import datetime
from flask import Flask, render_template, request, jsonify, send_file, redirect, url_for, flash, session
from flask_cors import CORS
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
import firebase_admin
from firebase_admin import credentials, auth as firebase_auth
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from dotenv import load_dotenv
import uuid
import io
import stripe  # Stripe payment integration

# CLIP services for semantic image matching
try:
    from services.clip_client import is_clip_available, get_text_embedding
    from services.image_matcher import pick_best_image_for_slide as clip_pick_best_image
    CLIP_IMPORT_SUCCESS = True
except ImportError as e:
    print(f"⚠️ CLIP services import failed: {e}")
    print("   → Install dependencies: pip install torch sentence-transformers")
    CLIP_IMPORT_SUCCESS = False
    clip_pick_best_image = None
    is_clip_available = lambda: False

TRANSLATION_CACHE = {}
CYRILLIC_RE = re.compile('[а-яА-Я]')

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)  # Enable CORS for cross-origin requests
app.secret_key = os.getenv('SECRET_KEY', 'your-secret-key-here-change-in-production')  # Needed for Flask-Login

# ============================================================================
# 🚨 CLIP INITIALIZATION - OPTIONAL FOR PRODUCTION
# ============================================================================
# For Railway deployment: CLIP dependencies are disabled to prevent crashes
# For local development: Install CLIP dependencies and enable

IS_RAILWAY = os.getenv('RAILWAY_ENVIRONMENT') is not None
CLIP_FORCE_DISABLE = os.getenv('CLIP_FORCE_DISABLE', 'false').lower() in ('true', '1', 'yes')

print("\n" + "="*70)
print("🔧 CLIP INITIALIZATION CHECK")
print("="*70)
print(f"Environment: {'Railway (Production)' if IS_RAILWAY else 'Local Development'}")
print(f"CLIP Force Disable: {CLIP_FORCE_DISABLE}")

if IS_RAILWAY or CLIP_FORCE_DISABLE:
    print("\n⚠️  CLIP DISABLED for this environment")
    print("   → Running in production mode without CLIP")
    print("   → Image search will use keyword-based matching only")
    print("="*70 + "\n")
    
    CLIP_AVAILABLE = False
    CLIP_IMPORT_SUCCESS = False
    
else:
    print("\n🔄 Attempting CLIP initialization...")
    print("   (This may take a minute on first run)\n")
    
    try:
        # STEP 1: Check PyTorch
        print("[1/5] Checking PyTorch...")
        import torch
        print(f"   ✅ PyTorch version: {torch.__version__}")
        print(f"   ✅ CUDA available: {torch.cuda.is_available()}")
        if torch.cuda.is_available():
            print(f"   ✅ CUDA version: {torch.version.cuda}")
            print(f"   ✅ GPU device: {torch.cuda.get_device_name(0)}")
            target_device = "CUDA"
        else:
            print(f"   ⚠️  CUDA not available, using CPU")
            target_device = "CPU"
        
        # STEP 2: Check CLIP library
        print("\n[2/5] Checking CLIP library...")
        import clip
        print(f"   ✅ CLIP library imported successfully")
        
        # STEP 3: Force load CLIP model
        print(f"\n[3/5] 🔥 Loading CLIP model (ViT-B/32 on {target_device})...")
        
        clip_load_start = time.perf_counter()
        device = "cuda" if torch.cuda.is_available() else "cpu"
        
        # Import clip_client and force initialization
        from services import clip_client
        
        # Force model load
        print(f"   → Loading model to {device.upper()}...")
        clip_client._device = device
        model, preprocess = clip.load("ViT-B/32", device=device)
        model.eval()
        
        # Set global variables
        clip_client._clip_model = model
        clip_client._clip_preprocess = preprocess
        clip_client._clip_available = True
        
        clip_load_time = time.perf_counter() - clip_load_start
        
        print(f"   ✅ CLIP model loaded successfully!")
        print(f"   ⏱️  Load time: {clip_load_time:.2f}s")
        
        # STEP 4: Load image cache
        print("\n[4/5] Loading image embedding cache...")
        clip_client._load_image_cache()
        cache_size = len(clip_client._image_embedding_cache)
        print(f"   ✅ Loaded {cache_size} cached embeddings")
        
        # STEP 5: Verify model is working
        print("\n[5/5] Testing CLIP functionality...")
        test_text = clip.tokenize(["test"]).to(device)
        with torch.no_grad():
            test_features = model.encode_text(test_text)
        print(f"   ✅ CLIP is functional (test embedding: {test_features.shape})")
        
        # SUCCESS!
        print("\n" + "="*70)
        print("🎯 CLIP INITIALIZATION COMPLETE - ALL SYSTEMS GO!")
        print("="*70)
        print(f"   🧠 Model: ViT-B/32")
        print(f"   💻 Device: {device.upper()}")
        print(f"   📊 Embedding dim: 512")
        print(f"   💾 Cached embeddings: {cache_size}")
        print(f"   ⏱️  Total init time: {clip_load_time:.2f}s")
        print("="*70 + "\n")
        
        # Set global flags
        CLIP_AVAILABLE = True
        CLIP_IMPORT_SUCCESS = True
        
    except ImportError as e:
        print("\n" + "="*70)
        print("⚠️  CLIP DEPENDENCIES NOT AVAILABLE")
        print("="*70)
        print(f"Import Error: {e}\n")
        print("→ Running without CLIP (keyword-based image search only)")
        print("→ For CLIP support, install: torch, torchvision, sentence-transformers")
        print("="*70 + "\n")
        
        CLIP_AVAILABLE = False
        CLIP_IMPORT_SUCCESS = False
        
    except Exception as e:
        print("\n" + "="*70)
        print("⚠️  CLIP INITIALIZATION FAILED")
        print("="*70)
        print(f"Error type: {type(e).__name__}")
        print(f"Error message: {e}\n")
        print("Full traceback:")
        print("─" * 70)
        import traceback
        traceback.print_exc()
        print("─" * 70)
        print("\n→ Running without CLIP (keyword-based image search only)")
        print("="*70 + "\n")
        
        CLIP_AVAILABLE = False
        CLIP_IMPORT_SUCCESS = False

# ============================================================================
# CLIP CONFIGURATION
# ============================================================================
CLIP_ENABLED = CLIP_AVAILABLE  # Based on initialization result
CLIP_SIMILARITY_THRESHOLD = float(os.getenv('CLIP_SIMILARITY_THRESHOLD', '0.30'))
CLIP_MIN_CANDIDATES = int(os.getenv('CLIP_MIN_CANDIDATES', '8'))
CLIP_MAX_CANDIDATES = int(os.getenv('CLIP_MAX_CANDIDATES', '20'))

print("\n" + "="*70)
print("🤖 CLIP CONFIGURATION")
print("="*70)
print(f"CLIP_ENABLED: {CLIP_ENABLED}")
print(f"CLIP_AVAILABLE: {CLIP_AVAILABLE}")
print(f"CLIP_SIMILARITY_THRESHOLD: {CLIP_SIMILARITY_THRESHOLD}")
print(f"CLIP_MIN_CANDIDATES: {CLIP_MIN_CANDIDATES}")
print(f"CLIP_MAX_CANDIDATES: {CLIP_MAX_CANDIDATES}")
if not CLIP_AVAILABLE:
    print("\n⚠️  Image search will use keyword-based matching only")
print("="*70 + "\n")

# ============================================================================
# TRANSLATION CONFIGURATION (Universal Layer)
# ============================================================================
# Universal translation layer for image search queries
# Supports multiple providers and can be enabled/disabled independently

# Main toggle for translation
TRANSLATION_ENABLED = os.getenv('TRANSLATION_ENABLED', 'false').lower() in ('true', '1', 'yes')

# Translation provider: 'none', 'libre', 'external'
TRANSLATION_PROVIDER = os.getenv('TRANSLATION_PROVIDER', 'none').lower()

# Target language for image searches (usually 'en' for better stock photo results)
TRANSLATION_TARGET_LANG = os.getenv('TRANSLATION_TARGET_LANG', 'en')

# LibreTranslate configuration (used when TRANSLATION_PROVIDER='libre')
LIBRETRANSLATE_URL = os.getenv('LIBRETRANSLATE_URL', 'http://localhost:5001')
LIBRETRANSLATE_TIMEOUT = int(os.getenv('LIBRETRANSLATE_TIMEOUT', '10'))

# External translation service configuration (used when TRANSLATION_PROVIDER='external')
EXTERNAL_TRANSLATE_URL = os.getenv('EXTERNAL_TRANSLATE_URL', '')
EXTERNAL_TRANSLATE_API_KEY = os.getenv('EXTERNAL_TRANSLATE_API_KEY', '')
EXTERNAL_TRANSLATE_TIMEOUT = float(os.getenv('EXTERNAL_TRANSLATE_TIMEOUT', '5.0'))

print("="*70)
print("🌐 TRANSLATION CONFIGURATION (Image Search)")
print("="*70)
print(f"TRANSLATION_ENABLED: {TRANSLATION_ENABLED}")
print(f"TRANSLATION_PROVIDER: {TRANSLATION_PROVIDER}")
print(f"TRANSLATION_TARGET_LANG: {TRANSLATION_TARGET_LANG}")

if not TRANSLATION_ENABLED:
    print("⚠️ Translation DISABLED for image search")
    print("   → Using original text for all image queries")
    print("   → Relying on CLIP semantic matching + multilingual photo stocks")
elif TRANSLATION_PROVIDER == 'none':
    print("ℹ️ Translation enabled but provider set to 'none'")
    print("   → No actual translation will occur")
    print("   → Using original text (same as TRANSLATION_ENABLED=false)")
elif TRANSLATION_PROVIDER == 'libre':
    print(f"✅ Translation provider: LibreTranslate")
    print(f"   → LibreTranslate URL: {LIBRETRANSLATE_URL}")
    print(f"   → Target language: {TRANSLATION_TARGET_LANG}")
    print("   → Note: Ensure LibreTranslate service is running")
elif TRANSLATION_PROVIDER == 'external':
    if EXTERNAL_TRANSLATE_URL:
        print(f"✅ Translation provider: External API")
        print(f"   → External URL: {EXTERNAL_TRANSLATE_URL}")
        print(f"   → Target language: {TRANSLATION_TARGET_LANG}")
        print(f"   → Timeout: {EXTERNAL_TRANSLATE_TIMEOUT}s")
    else:
        print("⚠️ Translation provider set to 'external' but EXTERNAL_TRANSLATE_URL not configured")
        print("   → Falling back to original text")
else:
    print(f"⚠️ Unknown translation provider: '{TRANSLATION_PROVIDER}'")
    print("   → Valid values: none, libre, external")
    print("   → Falling back to original text")

print("="*70 + "\n")

# ============================================================================
# IMAGE SEARCH MODE CONFIGURATION
# ============================================================================
# Control image search behavior: legacy (stable) vs advanced (experimental)

# USE_IMAGE_PROMPT: Whether to use LLM-generated image_prompt field
# - false (default): Legacy mode - ignores image_prompt, uses search_keyword/title/content
# - true: Advanced mode - uses image_prompt for better search queries
USE_IMAGE_PROMPT = os.getenv('USE_IMAGE_PROMPT', 'false').lower() in ('true', '1', 'yes')

# USE_STRICT_CLIP_FILTER: Whether CLIP should block images below threshold
# - false (default): Soft mode - CLIP only ranks, always picks best candidate
# - true: Strict mode - CLIP rejects images below CLIP_SIMILARITY_THRESHOLD
USE_STRICT_CLIP_FILTER = os.getenv('USE_STRICT_CLIP_FILTER', 'false').lower() in ('true', '1', 'yes')

print("="*70)
print("🖼️  IMAGE SEARCH MODE")
print("="*70)
print(f"USE_IMAGE_PROMPT: {USE_IMAGE_PROMPT}")
print(f"USE_STRICT_CLIP_FILTER: {USE_STRICT_CLIP_FILTER}")

if not USE_IMAGE_PROMPT and not USE_STRICT_CLIP_FILTER:
    print("📌 Mode: LEGACY (stable, production-ready)")
    print("   → Uses search_keyword/title/content for queries")
    print("   → CLIP only ranks candidates (no threshold blocking)")
    print("   → Maximum stability across RU/EN languages")
elif USE_IMAGE_PROMPT and not USE_STRICT_CLIP_FILTER:
    print("📌 Mode: ADVANCED with soft CLIP")
    print("   → Uses image_prompt when available")
    print("   → CLIP ranks but doesn't block images")
    print("   → Better quality with legacy fallback")
elif not USE_IMAGE_PROMPT and USE_STRICT_CLIP_FILTER:
    print("📌 Mode: LEGACY with strict CLIP")
    print("   → Uses search_keyword/title/content")
    print("   → CLIP can reject images below threshold")
    print("   → May skip images if relevance is low")
else:  # Both enabled
    print("📌 Mode: ADVANCED (experimental)")
    print("   → Uses image_prompt for queries")
    print("   → CLIP strictly filters by threshold")
    print("   → Best quality but may fail more often")
    print("   → Recommended only after thorough testing")

print("="*70 + "\n")

# ============================================================================
# DEVELOPMENT MODE CONFIGURATION
# ============================================================================
# Toggle payment verification for development/testing
# Set PAYMENTS_ENABLED=false in .env to disable payment checks
PAYMENTS_ENABLED = os.getenv('PAYMENTS_ENABLED', 'true').lower() in ('true', '1', 'yes')

if not PAYMENTS_ENABLED:
    print("⚠️  ========================================")
    print("⚠️  [DEV MODE] PAYMENTS DISABLED")
    print("⚠️  All payment checks will be bypassed")
    print("⚠️  This should ONLY be used for development/testing")
    print("⚠️  Set PAYMENTS_ENABLED=true for production")
    print("⚠️  ========================================")
else:
    print("✅ Payment verification: ENABLED (production mode)")

# ============================================================================
# Stripe Configuration
# ============================================================================
# Initialize Stripe with secret key from environment
STRIPE_SECRET_KEY = os.getenv('STRIPE_SECRET_KEY')
STRIPE_WEBHOOK_SECRET = os.getenv('STRIPE_WEBHOOK_SECRET')

if STRIPE_SECRET_KEY:
    stripe.api_key = STRIPE_SECRET_KEY
    print("✅ Stripe initialized successfully")
    print(f"   → API Key: {STRIPE_SECRET_KEY[:7]}...{STRIPE_SECRET_KEY[-4:]}")
else:
    print("⚠️ Stripe not configured: STRIPE_SECRET_KEY not found")
    print("   → Payment features will NOT work")

if STRIPE_WEBHOOK_SECRET:
    print("✅ Stripe webhook secret configured")
else:
    print("⚠️ STRIPE_WEBHOOK_SECRET not configured")
    print("   → Webhook signature verification will be skipped (INSECURE for production)")

# ============================================================================
# Firebase Admin SDK Initialization
# ============================================================================
# This section initializes Firebase Admin SDK for server-side authentication
# The SDK requires a service account key JSON file for secure communication
# with Firebase services

# Global flag to track Firebase initialization status
FIREBASE_INITIALIZED = False

try:
    # Step 1: Get service account key path from environment or use default
    firebase_cred_path = os.getenv('FIREBASE_SERVICE_ACCOUNT_KEY', 'serviceAccountKey.json')
    print(f"🔍 Checking Firebase service account key at: {firebase_cred_path}")
    
    # Step 2: Verify that the service account key file exists
    if not os.path.exists(firebase_cred_path):
        raise FileNotFoundError(f"Service account key not found at: {firebase_cred_path}")
    
    # Step 3: Validate JSON format and required fields with explicit UTF-8 encoding
    with open(firebase_cred_path, 'r', encoding='utf-8') as f:
        key_data = json.load(f)
        required_fields = ['type', 'project_id', 'private_key_id', 'private_key', 'client_email']
        missing_fields = [field for field in required_fields if field not in key_data]
        
        if missing_fields:
            raise ValueError(f"Service account key missing required fields: {missing_fields}")
        
        if key_data.get('type') != 'service_account':
            raise ValueError(f"Invalid key type: expected 'service_account', got '{key_data.get('type')}'")
        
        print(f"✅ Service account key validated: Project ID = {key_data.get('project_id')}")
    
    # Step 4: Initialize Firebase Admin SDK with credentials
    # Check if already initialized to prevent reinitialization errors
    if not firebase_admin._apps:
        cred = credentials.Certificate(firebase_cred_path)
        firebase_admin.initialize_app(cred)
        FIREBASE_INITIALIZED = True
        
        print("✅ Firebase Admin SDK initialized successfully")
        print("   → Token verification: ENABLED")
        print("   → User authentication: READY")
    else:
        print("ℹ️ Firebase Admin SDK already initialized")
        FIREBASE_INITIALIZED = True
    
except FileNotFoundError as e:
    print(f"⚠️ Firebase initialization failed: {e}")
    print("   → Firebase authentication will NOT work")
    print("   → Please add serviceAccountKey.json to project root")
    print("   → Download from: Firebase Console → Project Settings → Service Accounts")
    FIREBASE_INITIALIZED = False
    
except ValueError as e:
    print(f"⚠️ Firebase initialization failed: {e}")
    print("   → Invalid or corrupted service account key")
    print("   → Please download a new key from Firebase Console")
    FIREBASE_INITIALIZED = False
    
except Exception as e:
    print(f"⚠️ Firebase initialization error: {e}")
    print(f"   → Error type: {type(e).__name__}")
    print("   → Firebase authentication may not work correctly")
    FIREBASE_INITIALIZED = False

# Initialize Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'admin_login'
login_manager.login_message = 'Please log in to access this page.'

# Simple admin user storage
ADMIN_USERS = {
    'admin': {
        'password_hash': generate_password_hash(os.getenv('ADMIN_PASSWORD', 'admin123')),
        'id': 'admin'
    }
}

class User(UserMixin):
    def __init__(self, user_id, email=None, is_admin_user=False, name=None, picture=None):
        self.id = user_id
        self.email = email
        self.is_admin_user = is_admin_user
        self.name = name or email
        self.picture = picture

    def is_admin(self):
        return self.is_admin_user

# Lookup user by ID from SQLite
def get_user_by_id(user_id):
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute('SELECT id, email, name, picture, status FROM users WHERE id = ?', (user_id,))
        row = cursor.fetchone()
        conn.close()
        return dict(row) if row else None
    except Exception as e:
        print(f"Error fetching user by id: {e}")
        return None

@login_manager.user_loader
def load_user(user_id):
    # Admin shortcut
    if user_id in ADMIN_USERS:
        return User(user_id, is_admin_user=True)
    # Regular user
    try:
        int_id = int(user_id)
    except (ValueError, TypeError):
        return None
    user_data = get_user_by_id(int_id)
    if not user_data:
        return None
    return User(
        user_data['id'],
        email=user_data.get('email'),
        is_admin_user=False,
        name=user_data.get('name'),
        picture=user_data.get('picture')
    )

# Helper functions for Firebase user management
def get_or_create_firebase_user(firebase_uid, email, name='', picture=''):
    """
    Get existing user by Firebase UID or create new one
    Returns (user_data, error)
    """
    conn = None
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        # Try to find user by Firebase UID
        cursor.execute('SELECT * FROM users WHERE firebase_uid = ?', (firebase_uid,))
        user = cursor.fetchone()
        
        if user:
            # User exists, return their data
            return dict(user), None
        
        # Check if user exists by email (migrating from email/password auth)
        cursor.execute('SELECT * FROM users WHERE email = ?', (email,))
        user = cursor.fetchone()
        
        if user:
            # Update existing user with Firebase UID
            cursor.execute(
                'UPDATE users SET firebase_uid = ?, name = ?, picture = ? WHERE email = ?',
                (firebase_uid, name, picture, email)
            )
            conn.commit()
            cursor.execute('SELECT * FROM users WHERE email = ?', (email,))
            user = cursor.fetchone()
            return dict(user), None
        
        # Create new user
        cursor.execute(
            '''INSERT INTO users (email, firebase_uid, name, picture, status, free_credits)
               VALUES (?, ?, ?, ?, 'active', 3)''',
            (email, firebase_uid, name, picture)
        )
        conn.commit()
        user_id = cursor.lastrowid
        
        print(f"✅ New Firebase user created: {email} (ID: {user_id})")
        print(f"   → Free credits: 3 presentations")
        
        cursor.execute('SELECT * FROM users WHERE id = ?', (user_id,))
        user = cursor.fetchone()
        return dict(user), None
        
    except sqlite3.Error as e:
        if conn:
            conn.rollback()
        print(f"Database error in get_or_create_firebase_user: {e}")
        return None, f"Database error: {str(e)}"
    except Exception as e:
        if conn:
            conn.rollback()
        print(f"Error in get_or_create_firebase_user: {e}")
        return None, str(e)
    finally:
        if conn:
            conn.close()

# Helper functions for email/password authentication
def validate_email(email):
    """
    Validate email format
    Returns (is_valid, error_message)
    """
    if not email or len(email) < 3:
        return False, "Email is too short"
    if '@' not in email or '.' not in email:
        return False, "Invalid email format"
    if len(email) > 255:
        return False, "Email is too long"
    return True, None

def validate_password(password):
    """
    Validate password strength
    Returns (is_valid, error_message)
    """
    if not password or len(password) < 6:
        return False, "Password must be at least 6 characters"
    if len(password) > 128:
        return False, "Password is too long"
    return True, None

def create_user(email, password):
    """
    Create new user with email and password
    Returns (user_id, error)
    """
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # Check if user already exists
        cursor.execute('SELECT id FROM users WHERE email = ?', (email,))
        if cursor.fetchone():
            conn.close()
            return None, "User with this email already exists"
        
        # Create user
        password_hash = generate_password_hash(password)
        cursor.execute(
            '''INSERT INTO users (email, password_hash, status, free_credits)
               VALUES (?, ?, 'active', 3)''',
            (email, password_hash)
        )
        conn.commit()
        user_id = cursor.lastrowid
        conn.close()
        
        print(f"✅ New user created: {email} (ID: {user_id})")
        print(f"   → Free credits: 3 presentations")
        
        return user_id, None
        
    except Exception as e:
        print(f"Error creating user: {e}")
        return None, "Failed to create user"

def authenticate_user(email, password):
    """
    Authenticate user with email and password
    Returns (user_data, error)
    """
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        cursor.execute('SELECT * FROM users WHERE email = ?', (email,))
        user = cursor.fetchone()
        conn.close()
        
        if not user:
            return None, "Invalid email or password"
        
        if not user['password_hash']:
            return None, "Please sign in with Firebase/Google"
        
        if not check_password_hash(user['password_hash'], password):
            return None, "Invalid email or password"
        
        if user['status'] == 'blocked':
            return None, "Your account has been blocked"
        
        return dict(user), None
        
    except Exception as e:
        print(f"Error authenticating user: {e}")
        return None, "Authentication failed"

# Admin helper functions
def get_all_users():
    """
    Get all users from database
    """
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users ORDER BY registration_date DESC')
        users = [dict(row) for row in cursor.fetchall()]
        conn.close()
        return users
    except Exception as e:
        print(f"Error fetching users: {e}")
        return []

def update_user_status(user_id, status):
    """
    Update user status (active/blocked)
    """
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute('UPDATE users SET status = ? WHERE id = ?', (status, user_id))
        conn.commit()
        conn.close()
        return True
    except Exception as e:
        print(f"Error updating user status: {e}")
        return False

def delete_user(user_id):
    """
    Delete user and their presentations
    """
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        # Delete user's presentations first
        cursor.execute('DELETE FROM presentations WHERE user_id = ?', (user_id,))
        # Delete user
        cursor.execute('DELETE FROM users WHERE id = ?', (user_id,))
        conn.commit()
        conn.close()
        return True
    except Exception as e:
        print(f"Error deleting user: {e}")
        return False

# API Keys from environment variables
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
PEXELS_API_KEY = os.getenv('PEXELS_API_KEY')
UNSPLASH_ACCESS_KEY = os.getenv('UNSPLASH_ACCESS_KEY')  # Added Unsplash support

# ============================================================================
# IMAGE PROVIDER CONFIGURATION
# ============================================================================
# Configure image provider strategy: 'pexels', 'unsplash', or 'mixed'
# - 'pexels': Only use Pexels API
# - 'unsplash': Only use Unsplash API  
# - 'mixed': Try Pexels first, fallback to Unsplash (recommended)
IMAGE_PROVIDER_MODE = os.getenv('IMAGE_PROVIDER_MODE', 'mixed').lower()

# Configuration
OUTPUT_DIR = 'output'
IMAGE_CACHE_DIR = 'image_cache'

if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)
if not os.path.exists(IMAGE_CACHE_DIR):
    os.makedirs(IMAGE_CACHE_DIR)

# Initialize SQLite database for users
DB_PATH = 'users.db'

def init_db():
    """Initialize the database with users table"""
    conn = None
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # Create users table
        # Note: firebase_uid is added via migration below, not in CREATE TABLE
        # to avoid conflicts with existing databases
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                email TEXT UNIQUE NOT NULL,
                password_hash TEXT,
                google_id TEXT UNIQUE,
                name TEXT,
                picture TEXT,
                status TEXT DEFAULT 'active',
                registration_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # Create presentations table to track user activity
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS presentations (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                topic TEXT NOT NULL,
                num_slides INTEGER,
                filename TEXT,
                presentation_type TEXT DEFAULT 'business',
                creation_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id)
            )
        ''')
        
        # Create table to track used images (prevent duplicates)
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS used_images (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                image_url TEXT NOT NULL,
                image_query TEXT,
                used_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id)
            )
        ''')
        
        # Create index for faster image lookups
        cursor.execute("SELECT name FROM sqlite_master WHERE type='index' AND name='idx_used_images_user'")
        if not cursor.fetchone():
            try:
                cursor.execute('CREATE INDEX idx_used_images_user ON used_images(user_id, image_url)')
                print("✅ Migration: Created index on used_images table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: idx_used_images_user index may already exist - {e}")
        
        # Migration: Add missing columns to users table
        # Safe pattern: check column existence before adding to avoid errors
        cursor.execute("PRAGMA table_info(users)")
        existing_columns = [column[1] for column in cursor.fetchall()]
        
        # Add firebase_uid column if missing
        if 'firebase_uid' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN firebase_uid TEXT')
                print("✅ Migration: Added firebase_uid column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: firebase_uid column may already exist - {e}")
        
        # Add name column if missing
        if 'name' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN name TEXT')
                print("✅ Migration: Added name column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: name column may already exist - {e}")
        
        # Add picture column if missing
        if 'picture' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN picture TEXT')
                print("✅ Migration: Added picture column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: picture column may already exist - {e}")
        
        # Add google_id column if missing (for backward compatibility)
        if 'google_id' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN google_id TEXT')
                print("✅ Migration: Added google_id column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: google_id column may already exist - {e}")
        
        # Create unique index on firebase_uid if it doesn't exist
        cursor.execute("SELECT name FROM sqlite_master WHERE type='index' AND name='idx_firebase_uid'")
        if not cursor.fetchone():
            try:
                cursor.execute('CREATE UNIQUE INDEX idx_firebase_uid ON users(firebase_uid)')
                print("✅ Migration: Created unique index on firebase_uid column")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: idx_firebase_uid index may already exist - {e}")
        
        # Migration: Add presentation_type column if it doesn't exist
        cursor.execute("PRAGMA table_info(presentations)")
        columns = [column[1] for column in cursor.fetchall()]
        if 'presentation_type' not in columns:
            cursor.execute('ALTER TABLE presentations ADD COLUMN presentation_type TEXT DEFAULT "business"')
            print("✅ Migration: Added presentation_type column to presentations table")
        
        # Migration: Add Stripe-related columns to users table
        cursor.execute("PRAGMA table_info(users)")
        existing_columns = [column[1] for column in cursor.fetchall()]
        
        # Add stripe_customer_id column if missing
        if 'stripe_customer_id' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN stripe_customer_id TEXT')
                print("✅ Migration: Added stripe_customer_id column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: stripe_customer_id column may already exist - {e}")
        
        # Add subscription_plan column if missing
        if 'subscription_plan' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN subscription_plan TEXT DEFAULT "free"')
                print("✅ Migration: Added subscription_plan column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: subscription_plan column may already exist - {e}")
        
        # Add subscription_status column if missing
        if 'subscription_status' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN subscription_status TEXT DEFAULT "inactive"')
                print("✅ Migration: Added subscription_status column to users table")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: subscription_status column may already exist - {e}")
        
        # Migration: Add free_credits column to users table (3 free presentations for new users)
        if 'free_credits' not in existing_columns:
            try:
                cursor.execute('ALTER TABLE users ADD COLUMN free_credits INTEGER NOT NULL DEFAULT 3')
                print("✅ Migration: Added free_credits column to users table")
                print("   → New users will get 3 free presentations")
            except sqlite3.OperationalError as e:
                print(f"⚠️ Migration: free_credits column may already exist - {e}")
        
        conn.commit()
        
    except sqlite3.Error as e:
        print(f"❌ Database initialization error: {e}")
        if conn:
            conn.rollback()
        raise
    except Exception as e:
        print(f"❌ Unexpected error during database initialization: {e}")
        if conn:
            conn.rollback()
        raise
    finally:
        if conn:
            conn.close()

# Initialize database on startup
init_db()

# Presentation types configuration
# Presentation types configuration - REFACTORED TO 3 TYPES
PRESENTATION_TYPES = {
    'business': {
        'name_ru': 'Деловая презентация',
        'name_en': 'Business Presentation',
        'icon': '💼',
        'color': '#667eea',
        'temperature': 0.6,  # Confident, professional tone
        'structure': [
            {'title': 'Титульный слайд', 'description': 'Название, компания, контекст'},
            {'title': 'Контекст и проблема', 'description': 'Текущая ситуация, вызовы'},
            {'title': 'Наше решение/продукт', 'description': 'Предлагаемое решение'},
            {'title': 'Ценность и выгоды', 'description': 'Какую пользу приносит'},
            {'title': 'Ключевые функции', 'description': 'Основные возможности'},
            {'title': 'Результаты/кейсы', 'description': 'Достижения, примеры'},
            {'title': 'План/дорожная карта', 'description': 'Планы развития'},
            {'title': 'Команда/ресурсы', 'description': 'Кто реализует'},
            {'title': 'Следующие шаги/CTA', 'description': 'Призыв к действию'},
            {'title': 'Контакты', 'description': 'Как связаться'}
        ],
        'tips': 'Деловой уверенный тон без пафоса. Простой язык для бизнес-аудитории. Фокус на фактах и результатах.'
    },
    'scientific': {
        'name_ru': 'Научная презентация',
        'name_en': 'Scientific Presentation',
        'icon': '🔬',
        'color': '#27ae60',
        'temperature': 0.2,  # Academic, formal, highly detailed and precise
        'structure': [
            {'title': 'Титул и тема исследования', 'description': 'Название работы'},
            {'title': 'Введение и актуальность', 'description': 'Почему это важно'},
            {'title': 'Обзор литературы', 'description': 'Предыдущие работы'},
            {'title': 'Цель и задачи', 'description': 'Что исследуем'},
            {'title': 'Методология', 'description': 'Как исследовали'},
            {'title': 'Основные результаты', 'description': 'Данные и цифры'},
            {'title': 'Сравнение и обсуждение', 'description': 'Анализ результатов'},
            {'title': 'Выводы', 'description': 'Главные заключения'},
            {'title': 'Дальнейшие исследования', 'description': 'Перспективы'},
            {'title': 'Источники и благодарности', 'description': 'Литература'}
        ],
        'tips': 'Академический формальный стиль. Осторожные формулировки ("по данным исследований", "согласно литературе"). Максимум структурированности, минимум субъективности.'
    },
    'general': {
        'name_ru': 'Общая презентация',
        'name_en': 'General Presentation',
        'icon': '📊',
        'color': '#3498db',
        'temperature': 0.7,  # Friendly, explaining tone
        'structure': [
            {'title': 'Титульный слайд', 'description': 'Тема и цели'},
            {'title': 'Почему тема важна', 'description': 'Актуальность и значимость'},
            {'title': 'Ключевые понятия', 'description': 'Основные термины'},
            {'title': 'Основные идеи', 'description': 'Главные принципы'},
            {'title': 'Примеры из жизни', 'description': 'Практические кейсы'},
            {'title': 'Пошаговое объяснение', 'description': 'Детальный разбор'},
            {'title': 'Типичные ошибки', 'description': 'Чего избегать'},
            {'title': 'Краткое резюме', 'description': 'Основные выводы'},
            {'title': 'Вопросы для самопроверки', 'description': 'Проверка знаний'},
            {'title': 'Дополнительные ресурсы', 'description': 'Для углубленного изучения'}
        ],
        'tips': 'Дружелюбный объясняющий стиль. Много примеров и простых формулировок. Язык доступный для широкой аудитории.'
    }
}

# Supported languages
SUPPORTED_LANGUAGES = {
    'ru': 'Russian',
    'en': 'English',
    'es': 'Spanish',
    'zh': 'Chinese',
    'fr': 'French'
}

# AI role prompts per presentation type and language - REFACTORED TO 3 TYPES
def get_ai_role_prompt(presentation_type, language):
    """Get AI system role prompt based on presentation type and language"""
    prompts = {
        'business': {
            'ru': "Ты опытный бизнес-консультант и стратег. Создай деловую презентацию о компании, продукте или результатах. Используй деловой уверенный тон без пафоса, простой язык для бизнес-аудитории. Фокус на фактах, данных, конкретных результатах.",
            'en': "You are an experienced business consultant and strategist. Create a business presentation about company, product or results. Use confident professional tone without hype, simple language for business audience. Focus on facts, data, and concrete results.",
            'es': "Eres un consultor empresarial experimentado. Crea una presentación empresarial profesional en español con tono confiado y lenguaje simple.",
            'zh': "你是经验丰富的商业顾问。请用中文创建专业的商务演示文稿，使用自信的语调。",
            'fr': "Vous êtes un consultant en affaires expérimenté. Créez une présentation professionnelle en français avec un ton confiant."
        },
        'scientific': {
            'ru': "Ты научный исследователь с академическим опытом. Создай научную презентацию-доклад с фактами и цифрами. Используй академический формальный стиль, осторожные формулировки (\"по данным исследований\", \"в литературе описано\"). Максимум структурированности, минимум субъективности.",
            'en': "You are a scientific researcher with academic experience. Create a scientific presentation-report with facts and figures. Use academic formal style, careful formulations ('according to research', 'described in literature'). Maximum structure, minimum subjectivity.",
            'es': "Eres investigador científico. Crea una presentación científica en español con estilo formal y datos.",
            'zh': "你是科学研究员。请用中文创建科学演示文稿，使用正式风格。",
            'fr': "Vous êtes chercheur scientifique. Créez une présentation scientifique en français avec style formel."
        },
        'general': {
            'ru': "Ты профессиональный спикер и преподаватель. Создай общую презентацию для объяснения темы широкой аудитории. Используй дружелюбный объясняющий стиль, много примеров, простые формулировки. Доступный язык для школьников, студентов, любознательных людей.",
            'en': "You are a professional speaker and educator. Create a general presentation to explain a topic to broad audience. Use friendly explaining style, many examples, simple formulations. Accessible language for students and curious people.",
            'es': "Eres un educador profesional. Crea una presentación general en español con estilo amigable y muchos ejemplos.",
            'zh': "你是专业教师。请用中文创建通用演示文稿，使用友好的解释风格。",
            'fr': "Vous êtes un éducateur professionnel. Créez une présentation générale en français avec un style amical."
        }
    }
    # Default to business/en if not found
    return prompts.get(presentation_type, prompts['business']).get(language, prompts['business']['en'])

# System prompts per presentation type - REFACTORED TO 3 TYPES WITH BULLET-POINTS FOCUS
SYSTEM_PROMPTS = {
    'business': (
        'Ты опытный бизнес-консультант и аналитик с 10+ летним опытом. '
        'Создаёшь деловую презентацию о компании, продукте, результатах или бизнес-инициативе.\n\n'
        'ПРИНЦИПЫ:\n'
        '- Первый слайд - титульный: название темы, подзаголовок (для кого, о чём)\n'
        '- Проблема/контекст (почему это актуально)\n'
        '- Предлагаемое решение/продукт\n'
        '- Ценность и выгоды для клиента\n'
        '- Результаты, кейсы, метрики (цифры, данные)\n'
        '- Последний слайд: итоги + Call To Action (следующие шаги)\n\n'
        'ФОРМАТ ТЕКСТА:\n'
        '- На каждом слайде только тезисы (3–6 пунктов)\n'
        '- Каждый тезис - 1–2 предложения максимум\n'
        '- Никаких абзацев и длинных описаний\n'
        '- Каждый тезис несёт новую конкретную информацию\n\n'
        'СТИЛЬ: Деловой, конкретный, без художественных оборотов. Акцент на выгодах, результатах, цифрах, действиях.'
    ),
    'scientific': (
        'Ты учёный и исследователь. '
        'Создаёшь научную презентацию-доклад об исследовании, гипотезах, методах, результатах и выводах. '
        'Строго относишься к фактам.\n\n'
        'ПРИНЦИПЫ:\n'
        '- Первый слайд: название исследования + область\n'
        '- Введение: контекст, актуальность, обзор литературы\n'
        '- Цель и гипотезы исследования\n'
        '- Методология (кратко: как проводили исследование)\n'
        '- Основные результаты (данные, графики, таблицы)\n'
        '- Обсуждение и сравнение с литературой\n'
        '- Последний слайд: заключение (краткие выводы)\n\n'
        'ФОРМАТ ТЕКСТА:\n'
        '- На каждом слайде только тезисы (3–6 пунктов)\n'
        '- Каждый тезис - 1–2 предложения максимум\n'
        '- Разделяй факты, гипотезы и предположения\n'
        '- Используй осторожные формулировки: "по данным", "согласно литературе", "наблюдается"\n\n'
        'СТИЛЬ: Научный, формальный, структурированный. Максимум структуры, минимум субъективности.'
    ),
    'general': (
        'Ты опытный спикер и преподаватель. '
        'Создаёшь общую презентацию для объяснения сложных тем широкой аудитории. '
        'Умеешь просто и интересно рассказывать.\n\n'
        'ПРИНЦИПЫ:\n'
        '- Первый слайд: название темы + короткое описание\n'
        '- Почему тема важна (актуальность)\n'
        '- Основные понятия и термины (просто)\n'
        '- Ключевые идеи и практические шаги (по 1 идее на слайд)\n'
        '- Примеры, кейсы из жизни, типичные ошибки\n'
        '- Последний слайд: summary + что делать дальше (3–5 шагов)\n\n'
        'ФОРМАТ ТЕКСТА:\n'
        '- На каждом слайде только тезисы (3–6 пунктов)\n'
        '- Каждый тезис - 1–2 предложения максимум\n'
        '- Никаких абзацев, только чёткие пункты\n'
        '- Можно использовать вопросы к аудитории (умеренно)\n\n'
        'СТИЛЬ: Понятный, дружелюбный, с примерами из жизни. Доступно для школьников и студентов.'
    )
}

# Structure generator per type - REFACTORED TO 3 TYPES (5-10 SLIDES)
def get_slide_structure_by_type(presentation_type: str, num_slides: int):
    """
    Generate slide sequence for given presentation type.
    Returns list of slide roles/purposes based on type.
    Slides: 5-10 range (enforced).
    """
    seq = []
    t = presentation_type
    n = max(5, min(10, num_slides))  # Enforce 5-10 slides range
    
    if t == 'business':
        # Business: Title, Problem/Context, Solution, Value, Results, Plan, Team, CTA, Contacts
        seq = [
            'Title/Company',
            'Problem & Context',
            'Our Solution/Product',
            'Value & Benefits',
            'Key Features',
            'Results & Cases',
            'Plan/Roadmap',
            'Team/Resources',
            'Next Steps/CTA',
            'Contacts'
        ]
    elif t == 'scientific':
        # Scientific: Title, Intro, Literature, Goals, Methods, Results, Discussion, Conclusion, Future, References
        seq = [
            'Title & Research Topic',
            'Introduction & Relevance',
            'Literature Review',
            'Goals & Hypotheses',
            'Methodology',
            'Main Results',
            'Comparison & Discussion',
            'Conclusions',
            'Future Research',
            'References & Acknowledgments'
        ]
    else:  # 'general'
        # General: Title, Why Important, Key Concepts, Main Ideas, Examples, Explanation, Mistakes, Summary, Resources
        # REMOVED: 'Self-Check Questions' (quiz/assessment slides not allowed)
        seq = [
            'Title & Topic',
            'Why This Matters',
            'Key Concepts',
            'Main Ideas',
            'Real-Life Examples',
            'Step-by-Step Explanation',
            'Common Mistakes',
            'Summary',
            'Additional Resources'
        ]
    
    # Trim or expand to fit n slides
    if len(seq) >= n:
        return seq[:n]
    else:
        # Pad with last item if needed (rare case)
        return seq + [seq[-1]] * (n - len(seq))

# Get presentation type info safely
def get_presentation_type_info(presentation_type: str):
    return PRESENTATION_TYPES.get(presentation_type, PRESENTATION_TYPES['business'])

# Check if current user is admin
def is_admin():
    return current_user.is_authenticated and hasattr(current_user, 'is_admin_user') and current_user.is_admin_user


# ============================================================================
# UNIVERSAL TRANSLATION LAYER FOR IMAGE SEARCH
# ============================================================================

def external_translate(text: str, target_lang: str = 'en', source_lang: str = None) -> str:
    """
    Translate text using external HTTP API service.
    
    Universal template for external translation services (Google Translate API, DeepL, etc.)
    
    Args:
        text: Text to translate
        target_lang: Target language code (default: 'en')
        source_lang: Source language code (auto-detect if None)
    
    Returns:
        Translated text or original text if translation fails
    
    Configuration:
        EXTERNAL_TRANSLATE_URL: API endpoint
        EXTERNAL_TRANSLATE_API_KEY: API key (if required)
        EXTERNAL_TRANSLATE_TIMEOUT: Request timeout
    """
    if not EXTERNAL_TRANSLATE_URL:
        print(f"  ⚠️ External translation URL not configured, using original text")
        return text
    
    try:
        # Universal request template - adapt based on your provider
        # Example for Google Translate API, LibreTranslate, or similar
        headers = {}
        if EXTERNAL_TRANSLATE_API_KEY:
            headers['Authorization'] = f'Bearer {EXTERNAL_TRANSLATE_API_KEY}'
            # Or: headers['X-API-Key'] = EXTERNAL_TRANSLATE_API_KEY
        
        payload = {
            'q': text,
            'target': target_lang,
        }
        
        if source_lang:
            payload['source'] = source_lang
        
        print(f"  🌐 External translation: '{text[:40]}...' → {target_lang}")
        
        response = requests.post(
            EXTERNAL_TRANSLATE_URL,
            json=payload,
            headers=headers,
            timeout=EXTERNAL_TRANSLATE_TIMEOUT
        )
        
        if response.status_code == 200:
            data = response.json()
            # Adapt this based on response structure
            translated = data.get('translatedText') or data.get('translation') or data.get('text', '')
            translated = translated.strip()
            
            if translated:
                print(f"  ✅ External translation: '{text[:30]}' → '{translated[:30]}'")
                return translated
            else:
                print(f"  ⚠️ Empty translation response, using original")
                return text
        else:
            print(f"  ⚠️ External translation error {response.status_code}: {response.text[:100]}")
            return text
            
    except requests.exceptions.Timeout:
        print(f"  ⚠️ External translation timeout ({EXTERNAL_TRANSLATE_TIMEOUT}s), using original text")
        return text
    except requests.exceptions.ConnectionError as e:
        print(f"  ⚠️ External translation connection error: {e}")
        print(f"     Using original text")
        return text
    except Exception as e:
        print(f"  ⚠️ External translation exception: {e}")
        print(f"     Using original text")
        return text


def libre_translate(text: str, target_lang: str = 'en', source_lang: str = 'ru') -> str:
    """
    Translate text using LibreTranslate service.
    
    Args:
        text: Text to translate
        target_lang: Target language code (default: 'en')
        source_lang: Source language code (default: 'ru')
    
    Returns:
        Translated text or original text if translation fails
    """
    if not LIBRETRANSLATE_URL:
        print(f"  ⚠️ LibreTranslate URL not configured, using original text")
        return text
    
    try:
        payload = {
            'q': text,
            'source': source_lang,
            'target': target_lang
        }
        
        print(f"  🌐 LibreTranslate: '{text[:40]}...' → {target_lang} at {LIBRETRANSLATE_URL}")
        
        response = requests.post(
            f"{LIBRETRANSLATE_URL}/translate",
            json=payload,
            timeout=LIBRETRANSLATE_TIMEOUT
        )
        
        if response.status_code == 200:
            data = response.json()
            translated = data.get('translatedText', '').strip()
            # Sanitize minimal
            translated = re.sub(r'[^a-zA-Z\s]', '', translated)
            translated = ' '.join(translated.split())
            
            if translated:
                print(f"  ✅ LibreTranslate: '{text[:30]}' → '{translated[:30]}'")
                return translated
            else:
                print(f"  ⚠️ LibreTranslate returned empty, using original")
                return text
        else:
            print(f"  ⚠️ LibreTranslate error {response.status_code}: {response.text[:100]}")
            return text
            
    except requests.exceptions.Timeout:
        print(f"  ⚠️ LibreTranslate timeout ({LIBRETRANSLATE_TIMEOUT}s), using original text")
        return text
    except requests.exceptions.ConnectionError as e:
        print(f"  ⚠️ LibreTranslate connection error (service unavailable)")
        print(f"     Error: {e}")
        print(f"     Using original text")
        return text
    except Exception as e:
        print(f"  ⚠️ LibreTranslate exception: {e}")
        print(f"     Using original text")
        return text


def translate_for_image_search(text: str, source_lang: str = None, context: str = '') -> str:
    """
    Universal translation function for image search queries.
    
    This is the main entry point for all image search translations.
    Routes to appropriate provider based on configuration.
    
    Args:
        text: Text to translate (search query, keywords, etc.)
        source_lang: Source language code (auto-detected if None)
        context: Additional context (e.g., topic) for logging
    
    Returns:
        Translated text (or original if translation disabled/failed)
    
    Configuration:
        TRANSLATION_ENABLED: Master toggle
        TRANSLATION_PROVIDER: 'none', 'libre', 'external'
        TRANSLATION_TARGET_LANG: Target language (usually 'en')
    
    Examples:
        >>> translate_for_image_search("рост доходов")  # With CLIP: returns original
        >>> translate_for_image_search("рост доходов", source_lang='ru')  # Translates if enabled
    """
    if not text or not text.strip():
        return ''
    
    text = text.strip()
    
    # Auto-detect language if not specified
    if source_lang is None:
        source_lang = 'ru' if CYRILLIC_RE.search(text) else 'en'
    
    # Log context for debugging
    context_str = f" (context: {context})" if context else ""
    print(f"\n  🌐 Image search language: {source_lang}{context_str}")
    
    # Check if translation is disabled
    if not TRANSLATION_ENABLED:
        print(f"  ⚠️ Translation disabled (TRANSLATION_ENABLED=false)")
        print(f"     Using original query: '{text[:50]}...'")
        return text
    
    # Check if already in target language
    if source_lang == TRANSLATION_TARGET_LANG:
        print(f"  ℹ️ Text already in target language ({TRANSLATION_TARGET_LANG})")
        print(f"     Skipping translation: '{text[:50]}...'")
        return text
    
    # Check cache first
    cache_key = f"{context}|{text}".lower()
    if cache_key in TRANSLATION_CACHE:
        cached = TRANSLATION_CACHE[cache_key]
        print(f"  💾 From cache: '{text[:30]}' → '{cached[:30]}'")
        return cached
    
    print(f"  🌐 Translation: ENABLED, provider={TRANSLATION_PROVIDER}, target={TRANSLATION_TARGET_LANG}")
    
    # Route to appropriate provider
    translated = text  # Default to original
    
    if TRANSLATION_PROVIDER == 'none':
        print(f"  ℹ️ Provider set to 'none' - no translation")
        print(f"     Using original: '{text[:50]}...'")
        translated = text
        
    elif TRANSLATION_PROVIDER == 'libre':
        translated = libre_translate(text, TRANSLATION_TARGET_LANG, source_lang)
        
    elif TRANSLATION_PROVIDER == 'external':
        translated = external_translate(text, TRANSLATION_TARGET_LANG, source_lang)
        
    else:
        print(f"  ⚠️ Unknown provider '{TRANSLATION_PROVIDER}'")
        print(f"     Valid: 'none', 'libre', 'external'")
        print(f"     Using original: '{text[:50]}...'")
        translated = text
    
    # Cache the result
    if translated and translated != text:
        TRANSLATION_CACHE[cache_key] = translated
    
    return translated


# DEPRECATED: Legacy function for backward compatibility
def translate_keyword_to_english(keyword, topic=''):
    """
    DEPRECATED: Use translate_for_image_search() instead.
    
    Legacy wrapper for backward compatibility with existing code.
    Routes to new universal translation layer.
    """
    return translate_for_image_search(keyword, context=topic)


def detect_language(text):
    """
    Detect language: returns 'ru' if Cyrillic is present, else 'en'.
    """
    try:
        return 'ru' if CYRILLIC_RE.search(text or '') else 'en'
    except Exception:
        return 'en'


def detect_presentation_content_type(topic, slide_title, slide_content):
    """
    Detect the conceptual presentation type from content analysis.
    Returns one of: 'scientific', 'business', 'historical', 'technology', 
                    'philosophical', 'humanities', 'educational'
    
    This is different from user-selected presentation_type (business/scientific/general).
    This analyzes WHAT the content is about, not HOW it should be structured.
    """
    # Combine all text for analysis
    combined_text = f"{topic} {slide_title} {slide_content}".lower()
    
    # Scientific indicators (highest priority)
    scientific_keywords = [
        'research', 'study', 'experiment', 'hypothesis', 'data', 'methodology',
        'результаты', 'исследование', 'эксперимент', 'гипотеза', 'методология',
        'laboratory', 'лаборатория', 'scientific', 'науч', 'analysis', 'анализ',
        'theory', 'теория', 'conclusion', 'вывод', 'findings', 'evidence'
    ]
    
    # Technology indicators
    tech_keywords = [
        'software', 'algorithm', 'artificial intelligence', 'ai', 'machine learning',
        'программ', 'алгоритм', 'нейр', 'digital', 'цифров', 'computer', 'код',
        'blockchain', 'cloud', 'cybersecurity', 'кибербезопасность', 'innovation'
    ]
    
    # Business indicators
    business_keywords = [
        'market', 'revenue', 'profit', 'strategy', 'customer', 'product',
        'рынок', 'прибыль', 'стратегия', 'клиент', 'продукт', 'бизнес',
        'sales', 'продаж', 'investment', 'инвестиц', 'growth', 'рост',
        'company', 'компания', 'management', 'менеджмент'
    ]
    
    # Historical indicators
    historical_keywords = [
        'history', 'historical', 'century', 'век', 'историч', 'ancient',
        'medieval', 'средневеков', 'revolution', 'революц', 'war', 'войн',
        'empire', 'империя', 'dynasty', 'династия', 'civilization'
    ]
    
    # Philosophical/theoretical indicators
    philosophical_keywords = [
        'philosophy', 'филосо', 'concept', 'концеп', 'theory', 'теория',
        'ethics', 'этика', 'meaning', 'смысл', 'existence', 'сущест',
        'consciousness', 'сознание', 'logic', 'логика', 'metaphysics'
    ]
    
    # Humanities indicators (culture, art, society)
    humanities_keywords = [
        'culture', 'культур', 'art', 'искусство', 'society', 'общество',
        'literature', 'литератур', 'music', 'музык', 'painting', 'живопись',
        'social', 'социальн', 'anthropology', 'антропология', 'psychology'
    ]
    
    # Count matches for each category
    scores = {
        'scientific': sum(1 for kw in scientific_keywords if kw in combined_text),
        'technology': sum(1 for kw in tech_keywords if kw in combined_text),
        'business': sum(1 for kw in business_keywords if kw in combined_text),
        'historical': sum(1 for kw in historical_keywords if kw in combined_text),
        'philosophical': sum(1 for kw in philosophical_keywords if kw in combined_text),
        'humanities': sum(1 for kw in humanities_keywords if kw in combined_text)
    }
    
    # Get type with highest score (default to 'educational' if no clear match)
    max_score = max(scores.values())
    if max_score == 0:
        return 'educational'
    
    # Return the category with highest score
    detected_type = max(scores, key=scores.get)
    print(f"  🎯 Content type detected: {detected_type} (score: {max_score})")
    return detected_type


def generate_intelligent_image_query(slide_title, slide_content, topic, presentation_type, content_type=None):
    """
    Generate intelligent image search query based on:
    1. Presentation type (business/scientific/general) - user selected structure
    2. Content type (scientific/business/historical/etc) - detected from content
    3. Slide title and content keywords
    
    Returns: (english_query, original_language_query, image_type_category, description)
    
    Image type categories:
    - scientific: laboratory, research, diagrams, data visualization
    - corporate: team, office, graphs, business meeting
    - conceptual: abstract, infographic, diagram, visualization
    - historical: archival, portrait, historical scene, period-specific
    - tech: code, server, AI, digital, futuristic
    - real-world: people, nature, society, culture
    """
    # Auto-detect content type if not provided
    if content_type is None:
        content_type = detect_presentation_content_type(topic, slide_title, slide_content)
    
    # Extract keywords from title (2-3 main terms)
    title_words = re.findall(r'\b\w{4,}\b', slide_title.lower())  # Words 4+ chars
    
    # Extract keywords from first sentence of content (1-2 terms)
    first_sentence = slide_content.split('.')[0] if '.' in slide_content else slide_content[:100]
    content_words = re.findall(r'\b\w{5,}\b', first_sentence.lower())  # Words 5+ chars
    
    # Remove common stopwords
    stopwords = {
        'this', 'that', 'what', 'which', 'when', 'where', 'how', 'why',
        'это', 'этот', 'который', 'когда', 'где', 'как', 'почему',
        'introduction', 'conclusion', 'summary', 'overview',
        'введение', 'заключение', 'резюме', 'обзор'
    }
    
    title_keywords = [w for w in title_words if w not in stopwords][:3]
    content_keywords = [w for w in content_words if w not in stopwords][:2]
    
    # Determine image category and modifiers based on content type
    image_category = 'conceptual'  # default
    modifiers = []
    
    if content_type == 'scientific':
        image_category = 'scientific'
        modifiers = ['laboratory', 'research', 'scientific', 'experiment', 'data']
    elif content_type == 'business':
        image_category = 'corporate'
        modifiers = ['professional', 'business', 'modern office', 'team collaboration']
    elif content_type == 'technology':
        image_category = 'tech'
        modifiers = ['technology', 'digital', 'innovation', 'futuristic', 'code']
    elif content_type == 'historical':
        image_category = 'historical'
        modifiers = ['historical', 'archival', 'vintage', 'period', 'documentary']
    elif content_type == 'philosophical':
        image_category = 'conceptual'
        modifiers = ['abstract', 'concept', 'visualization', 'diagram', 'infographic']
    elif content_type == 'humanities':
        image_category = 'real-world'
        modifiers = ['people', 'culture', 'society', 'art', 'nature']
    else:  # educational or general
        image_category = 'conceptual'
        modifiers = ['educational', 'diagram', 'illustration', 'infographic']
    
    # Build search query components
    # Format: [main_keywords] + [category_modifier] + [quality_filter]
    
    # Translate keywords if needed
    translated_keywords = []
    for kw in (title_keywords + content_keywords):
        if CYRILLIC_RE.search(kw):
            translated = translate_keyword_to_english(kw, topic)
            if translated and translated != kw:
                translated_keywords.append(translated)
            else:
                translated_keywords.append(kw)
        else:
            translated_keywords.append(kw)
    
    # Select 1-2 best modifiers
    selected_modifiers = modifiers[:2]
    
    # Build final English query
    query_parts = translated_keywords[:2] + selected_modifiers[:1]
    english_query = ' '.join(query_parts)
    
    # Build original language query (for display)
    original_parts = (title_keywords + content_keywords)[:3]
    if detect_language(topic) == 'ru':
        original_query = ' '.join(original_parts)
    else:
        original_query = english_query
    
    # Generate description of what image should show
    if detect_language(topic) == 'ru':
        descriptions = {
            'scientific': 'Научное оборудование, лаборатория, исследователи за работой, научные диаграммы или графики',
            'corporate': 'Профессиональная рабочая среда, команда в офисе, деловая встреча, бизнес-графики',
            'tech': 'Современные технологии, компьютеры, код на экране, цифровые инновации, AI-системы',
            'historical': 'Исторические фотографии, архивные материалы, портреты исторических личностей',
            'conceptual': 'Абстрактная визуализация концепции, диаграмма идей, инфографика',
            'real-world': 'Реальные люди, культурные сцены, общество, природа, искусство'
        }
    else:
        descriptions = {
            'scientific': 'Scientific equipment, laboratory, researchers at work, scientific diagrams or charts',
            'corporate': 'Professional work environment, team in office, business meeting, business charts',
            'tech': 'Modern technology, computers, code on screen, digital innovations, AI systems',
            'historical': 'Historical photographs, archival materials, portraits of historical figures',
            'conceptual': 'Abstract concept visualization, idea diagrams, infographics',
            'real-world': 'Real people, cultural scenes, society, nature, art'
        }
    
    description = descriptions.get(image_category, descriptions['conceptual'])
    
    print(f"  🖼️ Image search: '{english_query}' | Category: {image_category}")
    
    return english_query, original_query, image_category, description


def generate_slide_content_in_language(topic, num_slides, language='en', presentation_type='business'):
    """
    Generate slide content using OpenAI ChatGPT API in the specified language
    with structure optimized for presentation type
    """
    try:
        print(f"Generating content in language: {language}, type: {presentation_type}")
        
        # Get presentation type info
        type_info = get_presentation_type_info(presentation_type)
        structure_guide = type_info.get('structure', [])
        tips = type_info.get('tips', '')
        temperature = type_info.get('temperature', 0.7)  # Get type-specific temperature
        
        # Build structure guidance string from type-specific sequence
        guided_sequence = get_slide_structure_by_type(presentation_type, num_slides)
        structure_text = "\n".join([f"- Slide {i+1}: {title}" for i, title in enumerate(guided_sequence)])
        
        headers = {
            'Authorization': f'Bearer {OPENAI_API_KEY}',
            'Content-Type': 'application/json'
        }
        
        # Get AI role prompt based on type and language
        language_name = SUPPORTED_LANGUAGES.get(language, 'English')
        system_prompt = get_ai_role_prompt(presentation_type, language)

        # Create prompt based on language and presentation type
        if language == 'ru':
            type_name_ru = type_info.get('name_ru', 'Презентация')
            prompt = f"""Создай структурированную презентацию на тему: "{topic}"
Количество слайдов: {num_slides}
Тип презентации: {type_name_ru}

РЕКОМЕНДУЕМАЯ СТРУКТУРА ДЛЯ ЭТОГО ТИПА:
{structure_text}

СОВЕТ ПО СТИЛЮ: {tips}

🎯 КРИТИЧЕСКИ ВАЖНЫЕ ТРЕБОВАНИЯ К КАЧЕСТВУ:

1. ГЛУБИНА И УНИКАЛЬНОСТЬ:
   • Каждый слайд должен содержать НЕОЖИДАННЫЕ факты, малоизвестные данные или оригинальные инсайты
   • НЕ ограничивайся общеизвестной информацией - углубляйся в нишевые детали
   • Каждый слайд уникален по углу рассмотрения темы
   • Раскрывай информацию для профессионалов и энтузиастов, а не для новичков

2. КОНКРЕТИКА И ДОКАЗАТЕЛЬСТВА:
   • ОБЯЗАТЕЛЬНО указывай: имена исследователей, годы исследований, названия институтов
   • Приводи точные цифры и статистику (не округляй: вместо "около 100" пиши "127 случаев")
   • Ссылайся на конкретные кейсы, практические примеры из реальной практики
   • Если данных нет - выскажи обоснованную гипотезу или критическое размышление

3. ЗАПРЕТ НА ШАБЛОНЫ:
   • СТРОГО ИЗБЕГАЙ фраз: "в современном мире", "в цифровую эпоху", "ключевой фактор", "новые возможности", "инновационные решения"
   • НЕ используй одинаковые структуры предложений на разных слайдах
   • НЕ заканчивай слайды похожими формулировками
   • Каждый слайд должен иметь свой авторский стиль изложения

4. ФОРМАТ КОНТЕНТА:
   • Каждый слайд: 3-6 тезисов (bullet points)
   • Каждый тезис - 1-2 предложения максимум
   • НЕ используй длинные абзацы - только чёткие пункты
   • Каждый пункт несёт новую конкретную информацию
   • Минимум 1-2 неожиданных факта на слайд
   • Используй аналогии, сравнения, критический анализ
   • Последний слайд: прогнозы, открытые вопросы, вызовы для будущих исследований

ПРИМЕР ГЛУБОКОГО КОНТЕНТА для темы "Нейросети для диагностики редких болезней у животных":
{{
  "slides": [
    {{
      "title": "Проблема гиподиагностики",
      "search_keyword": "veterinary diagnostics rare disease animals",
      "image_prompt": "veterinarian examining sick exotic pet in modern diagnostic clinic",
      "content": "Согласно исследованию Dr. Sarah Mitchell (Cornell University, 2022), только 12% редких заболеваний у домашних животных диагностируются при жизни. Основная причина — отсутствие у ветеринаров опыта распознавания атипичных симптомов. В случае синдрома Кушинга у хорьков средний срок до постановки диагноза составляет 8.3 месяца, что критично при средней продолжительности жизни 6-8 лет."
    }},
    {{
      "title": "Архитектура CNN для патологий",
      "search_keyword": "convolutional neural network medical imaging",
      "image_prompt": "medical imaging neural network analyzing microscopy pathology slides",
      "content": "Команда из UC Davis разработала сверточную сеть ResNet-152, обученную на 47,000 гистопатологических изображений экзотических животных. Точность детекции лимфомы у попугаев достигла 94.7%, превысив показатели опытных патологоанатомов (89.2%). Критический момент: сеть выявляет паттерны, невидимые человеческому глазу — анизоцитоз на уровне 3-5 микрон."
    }},
    {{
      "title": "Дилемма малых выборок",
      "search_keyword": "few shot learning medical AI",
      "image_prompt": "small dataset machine learning training process visualization",
      "content": "Для болезни фон Виллебранда у доберманов существует только 340 задокументированных случаев с подтвержденной биопсией. Техника few-shot learning с метрическими пространствами (Prototypical Networks) позволила достичь 78% точности при обучении всего на 15 примерах. Однако возникает риск переобучения: модель может запомнить артефакты конкретных клиник, а не истинные паттерны болезни."
    }},
    {{
      "title": "Открытые вызовы",
      "search_keyword": "AI challenges veterinary medicine future",
      "image_prompt": "diverse veterinary professionals discussing AI technology challenges",
      "content": "Три нерешенных вопроса тормозят внедрение: 1) Отсутствие стандартизации протоколов сбора данных между клиниками (89% баз данных несовместимы); 2) Этическая дилемма — кто несет ответственность при ошибке AI в диагнозе?; 3) Феномен 'distribution shift' — модели, обученные на данных из США, показывают падение точности на 23-31% при тестировании на азиатских породах. Требуются федеративные подходы к обучению."
    }}
  ]
}}

НЕПРАВИЛЬНО (шаблоны и общие фразы):
"Нейросети открывают новые возможности в ветеринарии. В современном мире технологии позволяют диагностировать болезни быстрее."

ПРАВИЛЬНО (конкретика и глубина):
"Алгоритм YOLO-v5, адаптированный группой Prof. Chen для рентгенограмм рептилий, обнаруживает метаболическую болезнь костей у игуан с чувствительностью 91.3% — на 34% выше, чем средний показатель герпетологов-практиков."

КРИТИЧЕСКИ ВАЖНО - СТРУКТУРА ОТВЕТА:
Для каждого слайда ОБЯЗАТЕЛЬНО верни:
- "title" - заголовок слайда
- "search_keyword" - ключевые слова для поиска картинок НА АНГЛИЙСКОМ (3-5 слов)
- "content" - содержимое слайда
- "image_prompt" - (ОПЦИОНАЛЬНО) подробное описание идеального изображения на АНГЛИЙСКОМ (5-12 слов)

ФОРМАТ ОТВЕТА:
Верни ТОЛЬКО валидный JSON в формате:
{{
  "slides": [
    {{"title": "...", "search_keyword": "...", "content": "..."}},
    ...
  ]
}}

Без markdown, без дополнительного текста."""
        elif language == 'es':
            prompt = f"""Crea una presentación estructurada sobre el tema: "{topic}"
Número de diapositivas: {num_slides}

IMPORTANTE: La presentación debe consistir en DECLARACIONES DE TESIS, no descripciones.

TESIS — una declaración clave que revela parte del tema.
NO solo describas, formula ideas y argumentos específicos.

ESTRUCTURA DE TESIS:
- Diapositiva 1: Idea principal del tema (declaración central)
- Diapositivas 2-{num_slides-1}: Aspectos clave, beneficios, aplicaciones
- Diapositiva {num_slides}: Conclusión, futuro, conclusión

Cada tesis debe:
✓ Ser una declaración específica directamente relacionada con "{topic}"
✓ Contener 2-3 oraciones precisas con DETALLES y EJEMPLOS CONCRETOS
✓ Desarrollar el tema principal
✓ Formar una cadena lógica con otras tesis
✓ EVITAR frases plantilla como "tecnología clave", "era digital", "sociedad moderna"
✓ Usar TERMINOLOGÍA ESPECÍFICA y hechos relevantes solo para "{topic}"

Para cada diapositiva, devuelve JSON con campos:
- "title": Título breve (2-3 palabras) específico para el tema
- "search_keyword": Palabras clave para búsqueda de imágenes en inglés (3-4 palabras)
- "content": TESIS — declaración específica (2-3 oraciones con detalles)

EJEMPLO de tesis correctas para "Perros":
{{
  "slides": [
    {{"title": "Evolución de los Perros", "search_keyword": "dog evolution wolf domestication", "content": "Los perros descienden de lobos aproximadamente hace 15,000 años a través de la domesticación. Las investigaciones genéticas muestran que los primeros perros aparecieron en Asia Oriental y se expandieron mundialmente con los humanos. Las razas modernas son resultado de la cría selectiva en los últimos 200 años."}},
    {{"title": "Razas y Funciones", "search_keyword": "dog breeds working dogs types", "content": "Existen más de 400 razas de perros reconocidas, cada una criada para tareas específicas. Las razas pastoriles (Border Collies, Pastores) manejan rebaños, las de caza (Retrievers, Spaniels) asisten en la caza, mientras que las razas guardianas (Dobermans, Rottweilers) protegen propiedades. Las razas de compañía (Chihuahuas, Terriers) se crían exclusivamente para compañía."}},
    {{"title": "Inteligencia Canina", "search_keyword": "dog intelligence training cognition", "content": "Los perros pueden memorizar hasta 165 palabras y gestos, comparable a las habilidades cognitivas de un niño de dos años. Los Border Collies se consideran la raza más inteligente, comprendiendo nuevos comandos tras solo 5 repeticiones. Las investigaciones muestran que los perros distinguen emociones humanas a través de expresiones faciales y tono de voz."}}
  ]
}}

INCORRECTO (frases plantilla):
"Los perros se están convirtiendo en un factor clave en la sociedad moderna. La adopción de estas tecnologías desbloquea nuevas posibilidades."

CORRECTO (hechos concretos):
"Los perros poseen un sentido del olfato 10,000 veces más agudo que los humanos debido a 300 millones de receptores olfativos. Esto les permite detectar drogas, explosivos e incluso diagnosticar cáncer en etapas tempranas."

Devuelve SOLO JSON válido sin texto adicional.

CRÍTICO: 
- Cada tesis debe contener HECHOS, NÚMEROS, EJEMPLOS relacionados con "{topic}"
- NO uses frases genéricas sobre "tecnología", "innovación", "futuro" sin especificaciones
- El título y contenido de cada diapositiva deben estar LÓGICAMENTE conectados
- Cada search_keyword debe ser DIFERENTE y específico"""
        elif language == 'zh':
            prompt = f"""创建关于主题 "{topic}" 的结构化演示文稿
幻灯片数量: {num_slides}

重要：演示文稿必须由论点陈述组成，而不是描述！

论点 — 揭示主题部分内容的关键陈述。
不要只是描述，要提出具体的想法和论据。

论点结构：
- 幻灯片 1: 主题的主要观点（核心陈述）
- 幻灯片 2-{num_slides-1}: 关键方面、优势、应用
- 幻灯片 {num_slides}: 结论、未来、要点

每个论点必须：
✓ 是与 "{topic}" 直接相关的具体陈述
✓ 包含 2-3 个带有具体细节和示例的精确句子
✓ 发展主要主题
✓ 与其他论点形成逻辑链
✓ 避免使用 "关键技术"、"数字时代"、"现代社会" 等模板短语
✓ 使用仅与 "{topic}" 相关的特定术语和事实

对于每张幻灯片，返回包含以下字段的 JSON：
- "title": 简短标题（2-3 个词），针对主题
- "search_keyword": 英文图像搜索关键词（3-4 个词）
- "content": 论点 — 具体陈述（2-3 个带细节的句子）

"狗" 的正确论点示例：
{{
  "slides": [
    {{"title": "狗的进化", "search_keyword": "dog evolution wolf domestication", "content": "狗大约在 15,000 年前通过驯化从狼进化而来。基因研究表明，第一批狗出现在东亚，并随着人类传播到世界各地。现代品种是过去 200 年选择性繁殖的结果。"}},
    {{"title": "品种和功能", "search_keyword": "dog breeds working dogs types", "content": "有超过 400 种被认可的狗品种，每种都为特定任务而培育。牧羊犬（边境牧羊犬、德国牧羊犬）管理牲畜，猎犬（寻回犬、西班牙猎犬）协助狩猎，而护卫犬（杜宾犬、罗威纳犬）保护财产。玩具犬（吉娃娃、梗犬）专门用于伴侣。"}},
    {{"title": "犬类智力", "search_keyword": "dog intelligence training cognition", "content": "狗能记住多达 165 个单词和手势，相当于两岁儿童的认知能力。边境牧羊犬被认为是最聪明的品种，只需 5 次重复就能理解新命令。研究表明狗能通过面部表情和语调区分人类情感。"}}
  ]
}}

错误（模板短语）：
"狗正在成为现代社会发展的关键因素。采用这些技术开启了新的可能性。"

正确（具体事实）：
"狗的嗅觉比人类敏锐 10,000 倍，因为它们拥有 3 亿个嗅觉受体。这使它们能够检测毒品、爆炸物，甚至在早期诊断癌症。"

仅返回有效的 JSON，不包含额外文本。

关键：
- 每个论点必须包含与 "{topic}" 相关的具体事实、数字、示例
- 不要使用没有具体说明的 "技术"、"创新"、"未来" 等通用短语
- 每张幻灯片的标题和内容必须在逻辑上相关联
- 每个 search_keyword 必须是不同的且具体的"""
        elif language == 'fr':
            prompt = f"""Créez une présentation structurée sur le sujet : "{topic}"
Nombre de diapositives : {num_slides}

IMPORTANT : La présentation doit consister en des DÉCLARATIONS DE THÈSE, pas des descriptions.

THÈSE — une déclaration clé qui révèle une partie du sujet.
Ne décrivez pas seulement, formulez des idées et arguments spécifiques.

STRUCTURE DES THÈSES :
- Diapositive 1 : Idée principale du sujet (déclaration centrale)
- Diapositives 2-{num_slides-1} : Aspects clés, avantages, applications
- Diapositive {num_slides} : Conclusion, avenir, point de vue

Chaque thèse doit :
✓ Être une déclaration spécifique directement liée à "{topic}"
✓ Contenir 2-3 phrases précises avec des DÉTAILS et EXEMPLES CONCRÉTS
✓ Développer le sujet principal
✓ Former une chaîne logique avec les autres thèses
✓ ÉVITER les phrases modèles comme "technologie clé", "ère numérique", "société moderne"
✓ Utiliser une TERMINOLOGIE SPÉCIFIQUE et des faits pertinents uniquement pour "{topic}"

Pour chaque diapositive, retournez JSON avec les champs :
- "title" : Titre bref (2-3 mots) spécifique au sujet
- "search_keyword" : Mots-clés pour recherche d'images en anglais (3-4 mots)
- "content" : THÈSE — déclaration spécifique (2-3 phrases avec détails)

EXEMPLE de thèses correctes pour "Chiens" :
{{
  "slides": [
    {{"title": "Évolution des Chiens", "search_keyword": "dog evolution wolf domestication", "content": "Les chiens descendent des loups il y a environ 15 000 ans par domestication. Les recherches génétiques montrent que les premiers chiens sont apparus en Asie de l'Est et se sont répandus dans le monde avec les humains. Les races modernes sont le résultat de l'élevage sélectif au cours des 200 dernières années."}},
    {{"title": "Races et Fonctions", "search_keyword": "dog breeds working dogs types", "content": "Plus de 400 races de chiens reconnues existent, chacune élevée pour des tâches spécifiques. Les races de berger (Border Collies, Bergers) gèrent les troupeaux, les races de chasse (Retrievers, Épagneuls) aident à la chasse, tandis que les races de garde (Dobermans, Rottweilers) protègent les propriétés. Les races de compagnie (Chihuahuas, Terriers) sont élevées exclusivement pour la compagnie."}},
    {{"title": "Intelligence Canine", "search_keyword": "dog intelligence training cognition", "content": "Les chiens peuvent mémoriser jusqu'à 165 mots et gestes, comparable aux capacités cognitives d'un enfant de deux ans. Les Border Collies sont considérés comme la race la plus intelligente, comprenant de nouvelles commandes après seulement 5 répétitions. Les recherches montrent que les chiens distinguent les émotions humaines par les expressions faciales et le ton de la voix."}}
  ]
}}

INCORRECT (phrases modèles) :
"Les chiens deviennent un facteur clé dans la société moderne. L'adoption de ces technologies débloque de nouvelles possibilités."

CORRECT (faits concrets) :
"Les chiens possèdent un sens de l'odorat 10 000 fois plus aigu que les humains grâce à 300 millions de récepteurs olfactifs. Cela leur permet de détecter des drogues, des explosifs et même de diagnostiquer le cancer à un stade précoce."

Retournez SEULEMENT du JSON valide sans texte supplémentaire.

CRITIQUE : 
- Chaque thèse doit contenir des FAITS CONCRÉTS, des NOMBRES, des EXEMPLES liés à "{topic}"
- N'utilisez PAS de phrases génériques sur "technologie", "innovation", "avenir" sans précisions
- Le titre et le contenu de chaque diapositive doivent être LIÉS LOGIQUEMENT
- Chaque search_keyword doit être DIFFÉRENT et spécifique"""
        else:  # Default to English
            type_name_en = type_info.get('name_en', 'Presentation')
            prompt = f"""Create a structured presentation on topic: "{topic}"
Number of slides: {num_slides}
Presentation type: {type_name_en}

RECOMMENDED STRUCTURE FOR THIS TYPE:
{structure_text}

STYLE ADVICE: {tips}

🎯 CRITICAL QUALITY REQUIREMENTS:

1. DEPTH AND UNIQUENESS:
   • Each slide must contain UNEXPECTED facts, little-known data, or original insights
   • DO NOT limit to common knowledge - dive into niche details
   • Each slide is unique in its angle of topic exploration
   • Target professionals and enthusiasts, not beginners

2. SPECIFICITY AND EVIDENCE:
   • MUST include: researcher names, study years, institution names
   • Provide exact numbers and statistics (don't round: instead of "about 100" write "127 cases")
   • Reference specific cases, practical examples from real practice
   • If no data available - state well-founded hypothesis or critical analysis

3. TEMPLATE BAN:
   • STRICTLY AVOID phrases: "in modern world", "in digital era", "key factor", "new opportunities", "innovative solutions"
   • DO NOT use identical sentence structures across slides
   • DO NOT end slides with similar formulations
   • Each slide must have its own authorial writing style

4. CONTENT FORMAT:
   • Each slide: 3-4 sentences with SPECIFIC details
   • Minimum 1-2 unexpected facts per slide
   • Use analogies, comparisons, critical analysis
   • Final slide: forecasts, open questions, challenges for future research

EXAMPLE OF DEEP CONTENT for "Neural networks for diagnosing rare animal diseases":
{{
  "slides": [
    {{
      "title": "Underdiagnosis Problem",
      "search_keyword": "veterinary diagnostics rare disease animals",
      "image_prompt": "veterinarian examining sick exotic pet in modern diagnostic clinic",
      "content": "According to Dr. Sarah Mitchell's study (Cornell University, 2022), only 12% of rare diseases in domestic animals are diagnosed during lifetime. Primary cause: veterinarians lack experience recognizing atypical symptoms. For Cushing's syndrome in ferrets, average time to diagnosis is 8.3 months, critical given 6-8 year lifespan."
    }},
    {{
      "title": "CNN Architecture for Pathology",
      "search_keyword": "convolutional neural network medical imaging",
      "image_prompt": "medical imaging neural network analyzing microscopy pathology slides",
      "content": "UC Davis team developed ResNet-152 convolutional network trained on 47,000 histopathological images of exotic animals. Lymphoma detection accuracy in parrots reached 94.7%, exceeding experienced pathologists (89.2%). Critical: network detects patterns invisible to human eye — anisocytosis at 3-5 micron level."
    }},
    {{
      "title": "Few-Shot Learning Dilemma",
      "search_keyword": "few shot learning medical AI",
      "image_prompt": "small dataset machine learning training process visualization",
      "content": "Only 340 documented biopsy-confirmed cases exist for von Willebrand disease in Dobermans. Few-shot learning with metric spaces (Prototypical Networks) achieved 78% accuracy training on just 15 examples. However, overfitting risk emerges: model may memorize artifacts of specific clinics rather than true disease patterns."
    }},
    {{
      "title": "Open Challenges",
      "search_keyword": "AI challenges veterinary medicine future",
      "image_prompt": "diverse veterinary professionals discussing AI technology challenges",
      "content": "Three unsolved issues hamper adoption: 1) Lack of standardized data collection protocols between clinics (89% databases incompatible); 2) Ethical dilemma — who bears responsibility for AI diagnostic errors?; 3) Distribution shift phenomenon — models trained on US data show 23-31% accuracy drop when tested on Asian breeds. Federated learning approaches required."
    }}
  ]
}}

INCORRECT (templates and generic phrases):
"Neural networks unlock new opportunities in veterinary medicine. Modern world technologies enable faster disease diagnosis."

CORRECT (specificity and depth):
"YOLO-v5 algorithm adapted by Prof. Chen's group for reptile X-rays detects metabolic bone disease in iguanas with 91.3% sensitivity — 34% higher than average herpetologist practitioners."

CRITICAL - RESPONSE STRUCTURE:
For each slide you MUST return:
- "title" - slide title
- "search_keyword" - keywords for image search IN ENGLISH (3-5 words)
- "content" - slide content
- "image_prompt" - (OPTIONAL) detailed description of ideal image in ENGLISH (5-12 words)

RESPONSE FORMAT:
Return ONLY valid JSON in format:
{{
  "slides": [
    {{"title": "...", "search_keyword": "...", "content": "..."}},
    ...
  ]
}}

No markdown, no additional text."""


        data = {
            'model': 'gpt-3.5-turbo',
            'messages': [
                {'role': 'system', 'content': f"{system_prompt}\n\nAlways respond with valid JSON only. Generate content in {language_name}."},
                {'role': 'user', 'content': prompt}
            ],
            'temperature': temperature,  # Use type-specific temperature (0.2 for scientific, 0.6 for business, 0.7 for general)
            'max_tokens': 2500  # Increased for detailed, in-depth responses
        }
        
        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers=headers,
            json=data,
            timeout=30
        )
        
        if response.status_code != 200:
            raise Exception(f"OpenAI API error: {response.status_code} - {response.text}")
        
        result = response.json()
        content = result['choices'][0]['message']['content'].strip()
        
        # Try to parse JSON from response
        # Remove markdown code blocks if present
        if content.startswith('```'):
            content = content.split('```')[1]
            if content.startswith('json'):
                content = content[4:]
            content = content.strip()
        
        slides_data = json.loads(content)
        return slides_data.get('slides', [])
        
    except json.JSONDecodeError as e:
        print(f"JSON parsing error: {e}")
        print(f"Response content: {content}")
        # Fail instead of generating low-quality fallback
        return None
    except Exception as e:
        print(f"Error generating content: {e}")
        # Fail instead of generating low-quality fallback
        return None


def create_fallback_slides(topic, num_slides, language='en'):
    """
    Create fallback slides if API fails with language support
    """
    slides = []
    
    # Language-specific fallback content
    if language == 'ru':
        slides.append({
            'title': f'{topic} изменяет мир',
            'search_keyword': f'{topic} innovation future technology',
            'content': f'{topic} становится ключевым фактором развития современного общества. Внедрение этих технологий открывает новые возможности для бизнеса и повседневной жизни. Понимание {topic} критически важно для успеха в цифровую эпоху.'
        })
        
        thesis_templates = [
            ('Ключевые преимущества', 'key benefits advantages', lambda t: f'{t} повышает эффективность работы и снижает издержки. Автоматизация процессов позволяет сосредоточиться на стратегических задачах. Компании, внедрившие {t}, получают конкурентное преимущество на рынке.'),
            ('Практическое применение', 'real world practical use', lambda t: f'Реальные кейсы показывают эффективность {t} в различных отраслях. От медицины до финансов, технология решает сложные задачи. Успешные примеры вдохновляют на дальнейшее внедрение.'),
            ('Вызовы и решения', 'challenges solutions problems', lambda t: f'Основные препятствия при внедрении {t} включают технические и организационные барьеры. Однако современные подходы позволяют эффективно преодолевать эти трудности. Правильная стратегия минимизирует риски и ускоряет адаптацию.'),
            ('Будущее технологии', 'future innovation development', lambda t: f'{t} будет играть всё более важную роль в ближайшие годы. Инвестиции в развитие этой области растут экспоненциально. Те, кто освоит {t} сегодня, станут лидерами завтрашнего дня.')
        ]
        
        for i in range(1, min(num_slides, len(thesis_templates) + 1)):
            title, keywords, content_func = thesis_templates[i - 1]
            slides.append({
                'title': title,
                'search_keyword': f'{topic} {keywords}',
                'content': content_func(topic)
            })
            
    elif language == 'es':
        slides.append({
            'title': f'{topic} Revoluciona',
            'search_keyword': f'{topic} innovacion futuro tecnologia',
            'content': f'{topic} está redefiniendo cómo abordamos los desafíos y oportunidades modernos. La adopción de estas tecnologías desbloquea nuevo potencial para negocios y vida diaria. Dominar {topic} es fundamental para el éxito en la era digital.'
        })
        
        thesis_templates = [
            ('Ventajas Clave', 'ventajas beneficios clave', lambda t: f'{t} mejora drásticamente la eficiencia mientras reduce costos operativos. La automatización permite a los equipos enfocarse en iniciativas estratégicas en lugar de tareas rutinarias. Las organizaciones que implementan {t} obtienen ventajas competitivas significativas en sus mercados.'),
            ('Impacto en el Mundo Real', 'impacto aplicaciones practicas', lambda t: f'Las historias de éxito demuestran la efectividad de {t} en diversas industrias. Desde la salud hasta las finanzas, la tecnología resuelve problemas anteriormente intratables. Estos ejemplos probados inspiran mayor adopción e innovación.'),
            ('Superando Desafíos', 'desafios soluciones problemas', lambda t: f'Los obstáculos principales para la adopción de {t} incluyen complejidad técnica y resistencia organizacional. Los marcos y metodologías modernos abordan efectivamente estas barreras. La planificación estratégica minimiza riesgos y acelera la implementación exitosa.'),
            ('Perspectiva Futura', 'futuro innovacion desarrollo', lambda t: f'{t} jugará un papel cada vez más vital en dar forma al mañana. La inversión en este campo crece exponencialmente año tras año. Los primeros adoptantes de {t} se posicionan como líderes del futuro.')
        ]
        
        for i in range(1, min(num_slides, len(thesis_templates) + 1)):
            title, keywords, content_func = thesis_templates[i - 1]
            slides.append({
                'title': title,
                'search_keyword': f'{topic} {keywords}',
                'content': content_func(topic)
            })
            
    elif language == 'zh':
        slides.append({
            'title': f'{topic} 革命',
            'search_keyword': f'{topic} innovation future technology',
            'content': f'{topic} 正在重塑我们应对现代挑战和机遇的方式。采用这些技术为业务和日常生活开启了新的可能性。掌握 {topic} 对于数字时代的成功至关重要。'
        })
        
        thesis_templates = [
            ('关键优势', 'key benefits advantages', lambda t: f'{t} 显著提高效率同时降低运营成本。自动化使团队能够专注于战略举措而非日常任务。实施 {t} 的组织在其市场中获得显著的竞争优势。'),
            ('现实世界影响', 'real world practical applications', lambda t: f'成功案例证明了 {t} 在不同行业的有效性。从医疗保健到金融，该技术解决了以前难以解决的问题。这些经过验证的例子激励着进一步的采用和创新。'),
            ('克服挑战', 'challenges solutions problems', lambda t: f'{t} 采用的主要障碍包括技术复杂性和组织阻力。现代框架和方法有效地解决了这些障碍。战略规划将风险降至最低并加速成功实施。'),
            ('未来展望', 'future innovation development', lambda t: f'{t} 将在塑造未来中发挥越来越重要的作用。该领域的投资正在逐年指数级增长。早期采用 {t} 的人将自己定位为未来的领导者。')
        ]
        
        for i in range(1, min(num_slides, len(thesis_templates) + 1)):
            title, keywords, content_func = thesis_templates[i - 1]
            slides.append({
                'title': title,
                'search_keyword': f'{topic} {keywords}',
                'content': content_func(topic)
            })
            
    elif language == 'fr':
        slides.append({
            'title': f'{topic} Révolution',
            'search_keyword': f'{topic} innovation future technologie',
            'content': f'{topic} redéfinit comment nous abordons les défis et opportunités modernes. L\'adoption de ces technologies débloque de nouvelles possibilités pour les entreprises et la vie quotidienne. Maîtriser {topic} est essentiel pour réussir à l\'ère numérique.'
        })
        
        thesis_templates = [
            ('Avantages Clés', 'avantages bénéfices clés', lambda t: f'{t} améliore drastiquement l\'efficacité tout en réduisant les coûts opérationnels. L\'automatisation permet aux équipes de se concentrer sur des initiatives stratégiques au lieu de tâches routinières. Les organisations implémentant {t} gagnent des avantages compétitifs significatifs sur leurs marchés.'),
            ('Impact Réel', 'impact applications pratiques', lambda t: f'Les histoires de réussite démontrent l\'efficacité de {t} dans diverses industries. De la santé aux finances, la technologie résout des problèmes auparavant intractables. Ces exemples éprouvés inspirent une adoption et une innovation supplémentaires.'),
            ('Surmonter les Défis', 'défis solutions problèmes', lambda t: f'Les obstacles principaux à l\'adoption de {t} incluent la complexité technique et la résistance organisationnelle. Les cadres et méthodologies modernes traitent efficacement ces barrières. La planification stratégique minimise les risques et accélère l\'implémentation réussie.'),
            ('Aperçu Futur', 'futur innovation développement', lambda t: f'{t} jouera un rôle de plus en plus vital dans façonner demain. L\'investissement dans ce domaine croît exponentiellement année après année. Les premiers adoptants de {t} se positionnent comme les leaders de l\'avenir.')
        ]
        
        for i in range(1, min(num_slides, len(thesis_templates) + 1)):
            title, keywords, content_func = thesis_templates[i - 1]
            slides.append({
                'title': title,
                'search_keyword': f'{topic} {keywords}',
                'content': content_func(topic)
            })
    else:  # Default to English
        slides.append({
            'title': f'{topic} Revolution',
            'search_keyword': f'{topic} innovation future technology',
            'content': f'{topic} is reshaping how we approach modern challenges and opportunities. The adoption of these technologies unlocks new potential for businesses and daily life. Mastering {topic} is critical for success in the digital age.'
        })
        
        thesis_templates = [
            ('Key Advantages', 'key benefits advantages', lambda t: f'{t} dramatically improves efficiency while reducing operational costs. Automation enables teams to focus on strategic initiatives instead of routine tasks. Organizations implementing {t} gain significant competitive advantages in their markets.'),
            ('Real-World Impact', 'real world practical applications', lambda t: f'Success stories demonstrate the effectiveness of {t} across diverse industries. From healthcare to finance, the technology solves previously intractable problems. These proven examples inspire further adoption and innovation.'),
            ('Overcoming Challenges', 'challenges solutions problems', lambda t: f'Primary obstacles to {t} adoption include technical complexity and organizational resistance. Modern frameworks and methodologies effectively address these barriers. Strategic planning minimizes risks and accelerates successful implementation.'),
            ('Future Outlook', 'future innovation development', lambda t: f'{t} will play an increasingly vital role in shaping tomorrow. Investment in this field is growing exponentially year over year. Early adopters of {t} position themselves as leaders of the future.')
        ]
        
        for i in range(1, min(num_slides, len(thesis_templates) + 1)):
            title, keywords, content_func = thesis_templates[i - 1]
            slides.append({
                'title': title,
                'search_keyword': f'{topic} {keywords}',
                'content': content_func(topic)
            })
    
    return slides


def get_cached_image_path(keywords):
    """
    Get cached image path based on keyword hash
    """
    cache_key = hashlib.md5(keywords.encode('utf-8')).hexdigest()
    cache_file = os.path.join(IMAGE_CACHE_DIR, f"{cache_key}.jpg")
    
    if os.path.exists(cache_file):
        print(f"  ⚡ Using cached image for '{keywords}'")
        return cache_file
    
    return None


def save_image_to_cache(image_data, keywords):
    """
    Save downloaded image to cache
    """
    try:
        cache_key = hashlib.md5(keywords.encode('utf-8')).hexdigest()
        cache_file = os.path.join(IMAGE_CACHE_DIR, f"{cache_key}.jpg")
        
        with open(cache_file, 'wb') as f:
            f.write(image_data.getvalue())
        
        return cache_file
    except Exception as e:
        print(f"  ⚠ Error caching image: {e}")
        return None


# ============================================================================
# Image Search API - Multi-source with fallback and rate limiting
# ============================================================================

# Rate limiting state
API_CALL_TIMES = {'pexels': [], 'unsplash': []}
MAX_CALLS_PER_MINUTE = {'pexels': 50, 'unsplash': 50}  # API limits

def can_make_api_call(service):
    """
    Check if we can make API call based on rate limits
    Returns True if allowed, False if rate limit exceeded
    """
    import time
    current_time = time.time()
    
    # Clean old calls (older than 60 seconds)
    API_CALL_TIMES[service] = [
        t for t in API_CALL_TIMES[service] 
        if current_time - t < 60
    ]
    
    # Check limit
    if len(API_CALL_TIMES[service]) >= MAX_CALLS_PER_MINUTE[service]:
        print(f"  ⚠ Rate limit reached for {service}")
        return False
    
    # Record this call
    API_CALL_TIMES[service].append(current_time)
    return True


# ============================================================================
# IMAGE PROVIDER LAYER - Multi-source image fetching
# ============================================================================
# This layer provides a unified interface for fetching images from multiple
# sources (Pexels, Unsplash) with automatic fallback and error handling

# ============================================================================
# USED IMAGES TRACKING SYSTEM
# ============================================================================
# Prevents duplicate images within presentations and across recent generations

def get_used_images_for_user(user_id, limit=100):
    """
    Get list of recently used image URLs for a user
    
    Args:
        user_id: User ID to fetch images for
        limit: Maximum number of recent images to return (default: 100)
    
    Returns:
        List of image URLs (strings)
    """
    if not user_id:
        return []
    
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute(
            '''SELECT image_url FROM used_images 
               WHERE user_id = ? 
               ORDER BY used_date DESC 
               LIMIT ?''',
            (user_id, limit)
        )
        rows = cursor.fetchall()
        conn.close()
        return [row[0] for row in rows]
    except Exception as e:
        print(f"⚠️ Error fetching used images: {e}")
        return []


def add_used_image(user_id, image_url, query=''):
    """
    Add an image to the used images tracking table
    
    Args:
        user_id: User ID who used the image
        image_url: URL of the image
        query: Search query used to find the image (optional)
    """
    if not user_id or not image_url:
        return
    
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute(
            '''INSERT INTO used_images (user_id, image_url, image_query)
               VALUES (?, ?, ?)''',
            (user_id, image_url, query)
        )
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"⚠️ Error adding used image: {e}")


def cleanup_old_used_images(user_id, keep_count=100):
    """
    Remove old used images beyond the keep_count limit
    Keeps the database from growing indefinitely
    
    Args:
        user_id: User ID to cleanup
        keep_count: Number of most recent images to keep (default: 100)
    """
    if not user_id:
        return
    
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # Delete all but the most recent keep_count images
        cursor.execute(
            '''DELETE FROM used_images 
               WHERE user_id = ? 
               AND id NOT IN (
                   SELECT id FROM used_images 
                   WHERE user_id = ? 
                   ORDER BY used_date DESC 
                   LIMIT ?
               )''',
            (user_id, user_id, keep_count)
        )
        
        deleted_count = cursor.rowcount
        conn.commit()
        conn.close()
        
        if deleted_count > 0:
            print(f"🧹 Cleaned up {deleted_count} old image entries for user {user_id}")
    except Exception as e:
        print(f"⚠️ Error cleaning up used images: {e}")

def fetch_images_from_pexels(query, count=1, retries=2):
    """
    Fetch images from Pexels API
    
    Args:
        query: Search query string
        count: Number of images to fetch (default: 1)
        retries: Number of retry attempts (default: 2)
    
    Returns:
        List of dicts with unified format:
        [
            {
                'url': 'https://...',
                'author': 'Photographer Name',
                'source': 'Pexels',
                'source_link': 'https://pexels.com/photo/...', 
                'attribution': 'Photo by Name on Pexels'
            }
        ]
        Returns empty list if no results or error
    """
    if not PEXELS_API_KEY:
        print("  ⚠ Pexels API key not configured")
        return []
    
    if not can_make_api_call('pexels'):
        return []
    
    for attempt in range(retries):
        try:
            query_clean = query.strip().lower()
            
            headers = {
                'Authorization': PEXELS_API_KEY
            }
            
            params = {
                'query': query_clean,
                'per_page': count,
                'orientation': 'landscape'
            }
            
            if attempt == 0:
                print(f"  → Pexels search: '{query_clean}'")
            
            response = requests.get(
                'https://api.pexels.com/v1/search',
                headers=headers,
                params=params,
                timeout=10
            )
            
            if response.status_code == 200:
                data = response.json()
                
                if data.get('photos') and len(data['photos']) > 0:
                    results = []
                    for photo in data['photos'][:count]:
                        results.append({
                            'url': photo['src']['large'],
                            'author': photo['photographer'],
                            'source': 'Pexels',
                            'source_link': photo.get('url', 'https://www.pexels.com'),
                            'attribution': f"Photo by {photo['photographer']} on Pexels"
                        })
                    print(f"  ✓ Pexels: Found {len(results)} image(s)")
                    return results
                else:
                    print(f"  ✗ No Pexels results for '{query_clean}'")
                    return []
            
            elif response.status_code == 429:  # Rate limit
                print(f"  ⚠ Pexels rate limit hit (attempt {attempt + 1}/{retries})")
                if attempt < retries - 1:
                    import time
                    time.sleep(1)  # Wait before retry
                    continue
                return []
            
            else:
                print(f"  ✗ Pexels API error: {response.status_code}")
                return []
        
        except requests.exceptions.Timeout:
            print(f"  ⚠ Pexels timeout (attempt {attempt + 1}/{retries})")
            if attempt < retries - 1:
                continue
            return []
        
        except Exception as e:
            print(f"  ✗ Pexels error: {e}")
            return []
    
    return []


def fetch_images_from_unsplash(query, count=1, retries=2):
    """
    Fetch images from Unsplash API
    
    Args:
        query: Search query string
        count: Number of images to fetch (default: 1)
        retries: Number of retry attempts (default: 2)
    
    Returns:
        List of dicts with unified format (same as fetch_images_from_pexels)
        Returns empty list if no results or error
    """
    if not UNSPLASH_ACCESS_KEY:
        return []  # Silent fail if not configured
    
    if not can_make_api_call('unsplash'):
        return []
    
    for attempt in range(retries):
        try:
            query_clean = query.strip().lower()
            
            headers = {
                'Authorization': f'Client-ID {UNSPLASH_ACCESS_KEY}'
            }
            
            params = {
                'query': query_clean,
                'per_page': count,
                'orientation': 'landscape'
            }
            
            if attempt == 0:
                print(f"  → Unsplash search: '{query_clean}'")
            
            response = requests.get(
                'https://api.unsplash.com/search/photos',
                headers=headers,
                params=params,
                timeout=10
            )
            
            if response.status_code == 200:
                data = response.json()
                
                if data.get('results') and len(data['results']) > 0:
                    results = []
                    for photo in data['results'][:count]:
                        results.append({
                            'url': photo['urls']['regular'],
                            'author': photo['user']['name'],
                            'source': 'Unsplash',
                            'source_link': photo['links']['html'],
                            'attribution': f"Photo by {photo['user']['name']} on Unsplash"
                        })
                    print(f"  ✓ Unsplash: Found {len(results)} image(s)")
                    return results
                else:
                    print(f"  ✗ No Unsplash results for '{query_clean}'")
                    return []
            
            elif response.status_code == 429:  # Rate limit
                print(f"  ⚠ Unsplash rate limit hit (attempt {attempt + 1}/{retries})")
                if attempt < retries - 1:
                    import time
                    time.sleep(1)
                    continue
                return []
            
            else:
                print(f"  ✗ Unsplash API error: {response.status_code}")
                return []
        
        except requests.exceptions.Timeout:
            print(f"  ⚠ Unsplash timeout (attempt {attempt + 1}/{retries})")
            if attempt < retries - 1:
                continue
            return []
        
        except Exception as e:
            print(f"  ✗ Unsplash error: {e}")
            return []
    
    return []


def get_images(query, count=1, mode=None):
    """
    Unified image fetching function with multi-source support
    
    This is the main function used by the presentation generator.
    It handles provider selection, fallback logic, and error handling.
    
    Args:
        query: Search query string
        count: Number of images to fetch (default: 1)
        mode: Override provider mode ('pexels', 'unsplash', 'mixed')
              If None, uses IMAGE_PROVIDER_MODE from config
    
    Returns:
        List of image dicts with unified format (url, author, source, etc.)
        Returns empty list if no images found
    
    Strategy:
        - 'pexels': Only try Pexels
        - 'unsplash': Only try Unsplash
        - 'mixed': Try Pexels first, fallback to Unsplash if needed
    """
    if mode is None:
        mode = IMAGE_PROVIDER_MODE
    
    results = []
    
    if mode == 'unsplash':
        # Unsplash only
        results = fetch_images_from_unsplash(query, count)
    
    elif mode == 'pexels':
        # Pexels only
        results = fetch_images_from_pexels(query, count)
    
    else:  # 'mixed' or default
        # Try Pexels first (primary source)
        results = fetch_images_from_pexels(query, count)
        
        if not results:
            # Fallback to Unsplash if Pexels failed or returned nothing
            print(f"  → Trying Unsplash as fallback...")
            results = fetch_images_from_unsplash(query, count)
    
    return results


def search_image(query):
    """
    Legacy wrapper for backward compatibility
    Searches for a single image and returns URL or None
    
    This function maintains compatibility with existing code that uses
    the old search_image() interface.
    """
    results = get_images(query, count=1)
    if results and len(results) > 0:
        return results[0]['url']
    return None


def search_pexels_image(query, retries=2):
    """
    Legacy wrapper for backward compatibility
    Returns: (image_url, attribution) or (None, None)
    """
    results = fetch_images_from_pexels(query, count=1, retries=retries)
    if results and len(results) > 0:
        img = results[0]
        return img['url'], img['attribution']
    return None, None


def search_unsplash_image(query, retries=2):
    """
    Legacy wrapper for backward compatibility  
    Returns: (image_url, attribution) or (None, None)
    """
    results = fetch_images_from_unsplash(query, count=1, retries=retries)
    if results and len(results) > 0:
        img = results[0]
        return img['url'], img['attribution']
    return None, None


def search_image_with_fallback(search_keyword, slide_title, main_topic, used_images, presentation_type='business', slide_content=''):
    """
    Search for image with intelligent query generation and multiple fallback attempts.
    Now uses AI-driven content analysis to select appropriate image types.
    
    Returns: (image_data, image_url, image_metadata) or (None, None, None)
    image_metadata: dict with 'query', 'category', 'description'
    """
    # Generate intelligent image search query
    english_query, original_query, image_category, description = generate_intelligent_image_query(
        slide_title=slide_title,
        slide_content=slide_content or '',
        topic=main_topic,
        presentation_type=presentation_type
    )
    
    attempts = []
    
    # Primary attempt: Intelligent query
    if english_query:
        attempts.append((english_query, f"Intelligent ({image_category})"))
    
    # Fallback 1: Original search keyword (if provided)
    if search_keyword and search_keyword.strip():
        if CYRILLIC_RE.search(search_keyword):
            translated = translate_keyword_to_english(search_keyword, main_topic)
            if translated and translated != english_query:
                attempts.append((translated, "Translated keyword"))
        elif search_keyword != english_query:
            attempts.append((search_keyword, "Original keyword"))
    
    # Fallback 2: Slide title
    if slide_title and slide_title != english_query:
        attempts.append((slide_title, "Slide title"))
    
    # Fallback 3: Main topic
    if main_topic and main_topic not in [a[0] for a in attempts[:3]]:
        attempts.append((main_topic, "Main topic"))
    
    metadata = {
        'query': english_query,
        'original_query': original_query,
        'category': image_category,
        'description': description
    }
    
    for query, attempt_name in attempts:
        if not query or query.strip() == "":
            continue
            
        print(f"  → Attempt: {attempt_name} - '{query}'")
        
        # Check cache first
        cached_path = get_cached_image_path(query)
        if cached_path and cached_path not in used_images:
            try:
                with open(cached_path, 'rb') as f:
                    image_data = io.BytesIO(f.read())
                return image_data, cached_path, metadata
            except:
                pass
        
        # Search on Pexels/Unsplash
        image_url = search_image(query)
        
        if image_url and image_url not in used_images:
            image_data = download_image(image_url)
            
            if image_data:
                # Save to cache
                cached_path = save_image_to_cache(image_data, query)
                return image_data, image_url, metadata
    
    print(f"  ✗ No unique image found after all attempts")
    return None, None, metadata


def search_image_for_slide(slide_title, slide_content, main_topic, exclude_images=None, presentation_type='business'):
    """
    MAIN ROUTING FUNCTION FOR IMAGE SEARCH
    
    Routes to either LEGACY or ADVANCED mode based on USE_IMAGE_PROMPT flag.
    This maintains backward compatibility while allowing opt-in to new features.
    
    - LEGACY mode (USE_IMAGE_PROMPT=false, default): Stable, simple keyword search
    - ADVANCED mode (USE_IMAGE_PROMPT=true): Uses image_prompt and enhanced pipeline
    
    Args:
        slide_title: Title of the slide
        slide_content: Main content/text of the slide
        main_topic: Overall presentation topic
        exclude_images: List of image URLs to exclude (previously used)
        presentation_type: Type of presentation (business/scientific/general)
    
    Returns:
        (image_data, image_url, query_used) or (None, None, None)
    """
    # Route based on USE_IMAGE_PROMPT flag
    if not USE_IMAGE_PROMPT:
        # LEGACY MODE: Ignore image_prompt, use search_keyword/title/content
        print(f"\n🏷️  MODE: LEGACY (USE_IMAGE_PROMPT=false)")
        return search_image_legacy_mode(
            slide_title=slide_title,
            slide_content=slide_content,
            main_topic=main_topic,
            exclude_images=exclude_images,
            presentation_type=presentation_type,
            search_keyword=None,  # Will be extracted from content
            language=None         # Will be auto-detected
        )
    else:
        # ADVANCED MODE: Use image_prompt if available
        print(f"\n🏷️  MODE: ADVANCED (USE_IMAGE_PROMPT=true)")
        return search_image_advanced_mode(
            slide_title=slide_title,
            slide_content=slide_content,
            main_topic=main_topic,
            exclude_images=exclude_images,
            presentation_type=presentation_type,
            image_prompt=None,  # Not available in this old signature
            language=None       # Will be auto-detected
        )


def search_image_in_curated_pool(clip_context_embedding, top_k: int = 5):
    """
    Search for images in curated pool using CLIP embeddings.
    
    STUB: Future implementation will use FAISS/vector database with pre-indexed curated images.
    
    Args:
        clip_context_embedding: CLIP embedding of slide context (numpy array)
        top_k: Number of top results to return
    
    Returns:
        List of image candidates (empty for now - stub implementation)
    
    Future:
        - Will maintain a curated pool of high-quality stock photos
        - Pre-computed CLIP embeddings stored in FAISS index
        - Fast vector similarity search
        - Metadata: tags, categories, license info
    """
    # STUB: Return empty list until curated pool is implemented
    return []


def search_image_legacy_mode(
    slide_title: str,
    slide_content: str,
    main_topic: str,
    exclude_images: list | None = None,
    presentation_type: str = 'business',
    search_keyword: str | None = None,
    language: str | None = None
):
    """
    LEGACY IMAGE SEARCH MODE - Maximum stability, minimal complexity
    
    This is the original, stable search behavior that works reliably across
    Russian and English presentations. It uses simple keyword-based search
    with optional CLIP ranking (soft mode - no threshold blocking).
    
    Key characteristics:
    - Uses search_keyword from LLM (or extracts from title/content)
    - Ignores image_prompt completely
    - CLIP only ranks candidates, never blocks images
    - Simple translation logic (if enabled)
    - Maximum compatibility and stability
    
    Args:
        slide_title: Title of the slide
        slide_content: Main content/text of the slide  
        main_topic: Overall presentation topic
        exclude_images: List of image URLs to exclude (previously used)
        presentation_type: Type of presentation (business/scientific/general)
        search_keyword: LLM-provided search keyword (preferred)
        language: Language of the slide content (auto-detected if None)
    
    Returns:
        (image_data, image_url, query_used) or (None, None, None)
    """
    if exclude_images is None:
        exclude_images = []
    
    print(f"\n🔍 [LEGACY] Searching image for slide: '{slide_title}'")
    
    # ========================================================================
    # STEP 1: Build search query from search_keyword or title/content
    # ========================================================================
    if search_keyword and search_keyword.strip():
        query = search_keyword.strip()
        print(f"  🎯 [LEGACY] Using search_keyword: '{query}'")
    else:
        # Extract keywords from title and content (old behavior)
        print(f"  ⚠️ [LEGACY] No search_keyword, extracting from title/content")
        
        text_for_query = f"{slide_title} {slide_content[:100]}"
        
        # Simple keyword extraction
        stopwords = {
            'the', 'is', 'at', 'which', 'on', 'a', 'an', 'as', 'are', 'was', 'were',
            'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did',
            'will', 'would', 'should', 'could', 'may', 'might', 'must',
            'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they',
            'это', 'этот', 'эта', 'эти', 'тот', 'та', 'те', 'и', 'в', 'на', 'по', 'с', 'у',
            'был', 'была', 'были', 'будет', 'будут', 'может', 'можно'
        }
        
        words = re.findall(r'\b\w{4,}\b', text_for_query.lower())
        keywords = [w for w in words if w not in stopwords][:5]
        
        if keywords:
            query = ' '.join(keywords[:3])  # Top 3 keywords
            print(f"  🎯 [LEGACY] Extracted keywords: {keywords[:3]}")
        else:
            query = slide_title
            print(f"  ⚠️ [LEGACY] No keywords, using title")
    
    # Auto-detect language if needed
    if language is None:
        language = detect_language(f"{slide_title} {slide_content[:50]}")
    
    print(f"  🌍 [LEGACY] Detected language: {language}")
    
    # ========================================================================
    # STEP 2: Apply translation if enabled
    # ========================================================================
    # Use universal translation layer
    query = translate_for_image_search(
        text=query,
        source_lang=language,
        context=f"legacy_search:{slide_title}"
    )
    
    print(f"  🔍 [LEGACY] Final search query: '{query}'")
    
    # ========================================================================
    # STEP 3: CLIP-enhanced search (SOFT MODE - no threshold blocking)
    # ========================================================================
    if CLIP_AVAILABLE:
        print(f"  🤖 [LEGACY] CLIP ranking: STRICT_FILTER={USE_STRICT_CLIP_FILTER}")
        
        # Fetch candidates (max 6 for speed optimization)
        candidate_count = 6
        candidates = get_images(query, count=candidate_count)
        
        if not candidates:
            print(f"  ⚠️ [LEGACY] No candidates for '{query}', trying title")
            candidates = get_images(slide_title, count=candidate_count)
        
        # Check minimum candidates threshold
        if candidates and len(candidates) < CLIP_MIN_CANDIDATES:
            print(f"  ⚠️ [LEGACY] Only {len(candidates)} candidates (< {CLIP_MIN_CANDIDATES} minimum)")
            print(f"     Skipping CLIP, using keyword search")
            candidates = []
        
        if candidates:
            print(f"  📊 [LEGACY] Found {len(candidates)} candidates, starting CLIP ranking...")
            print(f"     → CLIP Mode: {'STRICT (can reject)' if USE_STRICT_CLIP_FILTER else 'SOFT (ranks only)'}")
            print(f"     → Threshold: {CLIP_SIMILARITY_THRESHOLD if USE_STRICT_CLIP_FILTER else 0.0}")
            
            clip_start = time.perf_counter()
            
            # Build CLIP context (simple - no image_prompt)
            clip_context_text = f"{slide_title}. {slide_content[:60]}"
            print(f"  📋 [LEGACY] CLIP context: '{clip_context_text}...'")
            
            # Add description field if missing
            for candidate in candidates:
                if 'description' not in candidate:
                    candidate['description'] = (
                        candidate.get('attribution', '') or 
                        candidate.get('author', '') or 
                        slide_title
                    )
            
            # Use CLIP to rank images
            try:
                best_image = clip_pick_best_image(
                    slide_title=slide_title,
                    slide_content=slide_content,
                    image_candidates=candidates,
                    exclude_images=exclude_images,
                    similarity_threshold=CLIP_SIMILARITY_THRESHOLD if USE_STRICT_CLIP_FILTER else 0.0
                )
                
                clip_time = time.perf_counter() - clip_start
                print(f"  ⏱️  [LEGACY] CLIP processing completed in {clip_time:.2f}s")
                
                if best_image:
                    similarity = best_image.get('_clip_similarity', 'N/A')
                    source = best_image.get('source', 'Unknown')
                    
                    # In legacy mode, check if strict filter rejected image
                    if USE_STRICT_CLIP_FILTER and similarity != 'N/A' and similarity < CLIP_SIMILARITY_THRESHOLD:
                        print(f"  ❌ [LEGACY] CLIP rejected (similarity {similarity} < {CLIP_SIMILARITY_THRESHOLD})")
                        best_image = None
                    else:
                        image_url = best_image['url']
                        image_data = download_image(image_url)
                        
                        if image_data:
                            print(f"  ✅ [LEGACY] CLIP selected: {image_url[:50]}... (similarity={similarity}, source={source})")
                            return image_data, image_url, query
                        else:
                            print(f"  ⚠️ [LEGACY] Failed to download CLIP-selected image")
                else:
                    if USE_STRICT_CLIP_FILTER:
                        print(f"  ❌ [LEGACY] No image passed CLIP threshold ({CLIP_SIMILARITY_THRESHOLD})")
                    else:
                        print(f"  ⚠️ [LEGACY] CLIP ranking returned no result")
            except Exception as e:
                print(f"  ⚠️ [LEGACY] CLIP ranking failed: {e}")
        else:
            print(f"  ⚠️ [LEGACY] No candidates for CLIP ranking")
    else:
        # CLIP not available
        print(f"  ℹ️ [LEGACY] CLIP not available, using keyword search")
    
    # ========================================================================
    # STEP 4: Fallback to traditional keyword search
    # ========================================================================
    print(f"  🔍 [LEGACY] Fallback to keyword search")
    
    image_data, image_url, metadata = search_image_with_fallback(
        search_keyword=query,
        slide_title=slide_title,
        main_topic=main_topic,
        used_images=exclude_images,
        presentation_type=presentation_type,
        slide_content=slide_content
    )
    
    if image_url:
        print(f"  ✅ [LEGACY] Found image: {image_url[:60]}...")
        return image_data, image_url, query
    else:
        print(f"  ❌ [LEGACY] No suitable image found")
        return None, None, None


def build_image_search_query(
    slide_title: str,
    slide_content: str,
    image_prompt: str | None = None,
    language: str | None = None
) -> str:
    """
    Build optimal search query for image search based on available information.
    
    Priority:
    1. Use image_prompt if available (already in English, optimized for stock photos)
    2. Fall back to slide_title + content keywords
    3. Apply translation if needed (based on TRANSLATION_ENABLED/PROVIDER)
    
    Args:
        slide_title: Title of the slide
        slide_content: Content of the slide
        image_prompt: LLM-generated image description in English (preferred)
        language: Language of the slide (auto-detected if None)
    
    Returns:
        Search query string optimized for Pexels/Unsplash
    
    Examples:
        >>> build_image_search_query(
        ...     "Market Analysis",
        ...     "Our revenue grew...",
        ...     "business team analyzing financial charts in modern office"
        ... )
        "business team analyzing financial charts in modern office"
        
        >>> build_image_search_query(
        ...     "Анализ рынка",
        ...     "Наши доходы выросли...",
        ...     None,
        ...     language='ru'
        ... )
        # Returns translated query or original based on TRANSLATION_ENABLED
    """
    # ========================================================================
    # PRIORITY 1: Use image_prompt if available (best option)
    # ========================================================================
    if image_prompt and image_prompt.strip():
        query = image_prompt.strip()
        print(f"  🖼️ Image prompt: '{query}'")
        # image_prompt should already be in English, optimized for stock photos
        # No translation needed
        return query
    
    # ========================================================================
    # PRIORITY 2: Build query from title + content
    # ========================================================================
    print(f"  ⚠️ No image_prompt provided, building from title/content")
    
    # Combine title and short content snippet
    text_for_query = f"{slide_title} {slide_content[:100]}"
    
    # Auto-detect language if not specified
    if language is None:
        language = detect_language(text_for_query)
    
    print(f"  🌐 Detected language: {language}")
    
    # Extract keywords (simplified - reuse existing logic)
    stopwords = {
        'the', 'is', 'at', 'which', 'on', 'a', 'an', 'as', 'are', 'was', 'were',
        'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did',
        'will', 'would', 'should', 'could', 'may', 'might', 'must',
        'this', 'that', 'these', 'those', 'i', 'you', 'he', 'she', 'it', 'we', 'they',
        'это', 'этот', 'эта', 'эти', 'тот', 'та', 'те', 'и', 'в', 'на', 'по', 'с', 'у',
        'был', 'была', 'были', 'будет', 'будут', 'может', 'можно'
    }
    
    words = re.findall(r'\b\w{4,}\b', text_for_query.lower())
    keywords = [w for w in words if w not in stopwords][:5]
    
    if keywords:
        query = ' '.join(keywords[:3])  # Top 3 keywords
        print(f"  🎯 Extracted keywords: {keywords[:3]}")
    else:
        query = slide_title
        print(f"  ⚠️ No keywords extracted, using title")
    
    # ========================================================================
    # TRANSLATION: Use universal translation layer
    # ========================================================================
    # Translate query if:
    # - TRANSLATION_ENABLED=true
    # - TRANSLATION_PROVIDER != 'none'
    # - language is not already target language
    
    query = translate_for_image_search(
        text=query,
        source_lang=language,
        context=f"image_search:{slide_title}"
    )
    
    print(f"  🔍 Final search query: '{query}'")
    return query


def search_image_advanced_mode(
    slide_title: str,
    slide_content: str,
    main_topic: str,
    exclude_images: list | None = None,
    presentation_type: str = 'business',
    image_prompt: str | None = None,
    language: str | None = None
):
    """
    ADVANCED IMAGE SEARCH MODE - Uses image_prompt and enhanced pipeline
    
    This mode uses the newer, more sophisticated search pipeline with:
    - image_prompt from LLM for better search queries
    - Universal translation layer
    - CLIP semantic matching with configurable threshold filtering
    - Curated pool support (stub for future)
    
    Behavior depends on USE_STRICT_CLIP_FILTER:
    - false: CLIP ranks but never blocks (soft mode)
    - true: CLIP can reject images below threshold (strict mode)
    
    Args:
        slide_title: Title of the slide
        slide_content: Main content/text of the slide
        main_topic: Overall presentation topic
        exclude_images: List of image URLs to exclude (previously used)
        presentation_type: Type of presentation (business/scientific/general)
        image_prompt: LLM-generated image description in English (NEW)
        language: Language of the slide content (auto-detected if None)
    
    Returns:
        (image_data, image_url, query_used) or (None, None, None)
    """
    if exclude_images is None:
        exclude_images = []
    
    print(f"\n🔍 [ADVANCED] Searching image for slide: '{slide_title}'")
    print(f"  🔧 [ADVANCED] STRICT_FILTER={USE_STRICT_CLIP_FILTER}")
    
    # ========================================================================
    # NEW: Build search query using image_prompt or fallback
    # ========================================================================
    search_query = build_image_search_query(
        slide_title=slide_title,
        slide_content=slide_content,
        image_prompt=image_prompt,
        language=language
    )
    
    # ========================================================================
    # FUTURE: Try curated pool first (stub for now)
    # ========================================================================
    if CLIP_AVAILABLE and image_prompt:
        # Get CLIP embedding for slide context
        try:
            from services.clip_client import get_text_embedding
            clip_context = f"{slide_title}. {slide_content[:100]}. {image_prompt or ''}"
            context_embedding = get_text_embedding(clip_context)
            
            # Try curated pool (returns empty list for now - stub)
            curated_candidates = search_image_in_curated_pool(context_embedding, top_k=5)
            
            if curated_candidates:
                print(f"  🌟 [ADVANCED] Found {len(curated_candidates)} images in curated pool")
                # TODO: Implement curated pool ranking and selection
                # For now, falls through to regular search
        except Exception as e:
            print(f"  ⚠️ [ADVANCED] Curated pool search failed: {e}")
    
    # ========================================================================
    # CLIP-ENHANCED IMAGE SEARCH (from Pexels/Unsplash)
    # ========================================================================
    if CLIP_AVAILABLE:
        print(f"  🤖 [ADVANCED] Using CLIP semantic matching")
        print(f"     Threshold: {CLIP_SIMILARITY_THRESHOLD}, Min candidates: {CLIP_MIN_CANDIDATES}")
        
        # Determine candidate count (max 6 for speed optimization)
        candidate_count = 6
        
        # Fetch candidates using the built search query
        candidates = get_images(search_query, count=candidate_count)
        
        if not candidates:
            print(f"  ⚠️ [ADVANCED] No candidates for '{search_query}', trying title")
            # Try with slide title as fallback
            candidates = get_images(slide_title, count=candidate_count)
        
        # Check if we have minimum required candidates
        if candidates and len(candidates) < CLIP_MIN_CANDIDATES:
            print(f"  ⚠️ [ADVANCED] Only {len(candidates)} candidates (< {CLIP_MIN_CANDIDATES} minimum)")
            print(f"     Skipping CLIP, falling back to keyword search")
            candidates = []  # Force fallback
        
        if candidates:
            print(f"  📊 [ADVANCED] Found {len(candidates)} candidates, applying CLIP ranking...")
            
            # Enhanced CLIP context with image_prompt
            if image_prompt:
                clip_context_text = f"{slide_title}. {slide_content[:60]}. Target: {image_prompt}"
            else:
                clip_context_text = f"{slide_title}. {slide_content[:60]}"
            
            print(f"  📝 CLIP context: '{clip_context_text}...'")
            
            # Add description field if missing
            for candidate in candidates:
                if 'description' not in candidate:
                    candidate['description'] = (
                        candidate.get('attribution', '') or 
                        candidate.get('author', '') or 
                        slide_title
                    )
            
            # Use CLIP to pick best matching image
            # In soft mode (USE_STRICT_CLIP_FILTER=false), pass threshold=0.0 to never block
            # In strict mode (USE_STRICT_CLIP_FILTER=true), use actual threshold
            effective_threshold = CLIP_SIMILARITY_THRESHOLD if USE_STRICT_CLIP_FILTER else 0.0
            
            try:
                best_image = clip_pick_best_image(
                    slide_title=slide_title,
                    slide_content=slide_content + (f" Image target: {image_prompt}" if image_prompt else ""),
                    image_candidates=candidates,
                    exclude_images=exclude_images,
                    similarity_threshold=effective_threshold
                )
                
                if best_image:
                    similarity = best_image.get('_clip_similarity', 'N/A')
                    source = best_image.get('source', 'Unknown')
                    
                    # Check if strict mode rejected the image
                    if USE_STRICT_CLIP_FILTER and similarity != 'N/A' and similarity < CLIP_SIMILARITY_THRESHOLD:
                        print(f"  ❌ [ADVANCED] CLIP rejected (similarity {similarity} < {CLIP_SIMILARITY_THRESHOLD})")
                        print(f"     Reason: Strict filter enabled, threshold not met")
                    else:
                        image_url = best_image['url']
                        image_data = download_image(image_url)
                        
                        if image_data:
                            mode_suffix = "(strict)" if USE_STRICT_CLIP_FILTER else "(soft)"
                            print(f"  ✅ [ADVANCED] CLIP selected {mode_suffix}: {image_url[:50]}...")
                            print(f"     similarity={similarity}, source={source}")
                            return image_data, image_url, search_query
                        else:
                            print(f"  ⚠️ [ADVANCED] Failed to download CLIP-selected image")
                else:
                    if USE_STRICT_CLIP_FILTER:
                        print(f"  ❌ [ADVANCED] No image passed CLIP threshold ({CLIP_SIMILARITY_THRESHOLD})")
                    else:
                        print(f"  ⚠️ [ADVANCED] CLIP ranking returned no result")
            except Exception as e:
                print(f"  ⚠️ [ADVANCED] CLIP ranking failed: {e}")
        else:
            print(f"  ⚠️ [ADVANCED] No candidates for CLIP ranking")
    else:
        # CLIP not available
        if CLIP_ENABLED:
            print(f"  ⚠️ [ADVANCED] CLIP enabled but not available (initialization failed)")
        else:
            print(f"  ℹ️ [ADVANCED] CLIP disabled (CLIP_ENABLED=false)")
        print(f"     Using keyword search only")
    
    # ========================================================================
    # FALLBACK: Traditional keyword-based search
    # ========================================================================
    print(f"  🔍 [ADVANCED] Fallback to keyword search")
    
    # Use existing intelligent search with duplicate prevention
    image_data, image_url, metadata = search_image_with_fallback(
        search_keyword=search_query,
        slide_title=slide_title,
        main_topic=main_topic,
        used_images=exclude_images,
        presentation_type=presentation_type,
        slide_content=slide_content
    )
    
    if image_url:
        print(f"  ✅ [ADVANCED] Found image: {image_url[:60]}...")
        return image_data, image_url, search_query
    else:
        print(f"  ❌ [ADVANCED] No suitable image found")
        return None, None, None


# Backward compatibility alias
search_image_for_slide_enhanced = search_image_advanced_mode


def is_libretranslate_available():
    """
    Check if LibreTranslate service is available.
    Returns False immediately if translation is disabled or provider is not 'libre'.
    """
    try:
        if not TRANSLATION_ENABLED:
            return False
        if TRANSLATION_PROVIDER != 'libre':
            return False
        if not LIBRETRANSLATE_URL:
            return False
        resp = requests.get(f"{LIBRETRANSLATE_URL}/languages", timeout=LIBRETRANSLATE_TIMEOUT)
        return resp.status_code == 200
    except Exception:
        return False


def download_image(url):
    """
    Download image from URL and return as bytes
    Security: Limit image size to prevent memory issues
    """
    MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10 MB limit
    
    try:
        response = requests.get(url, timeout=10, stream=True)
        if response.status_code == 200:
            # Check content length if available
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > MAX_IMAGE_SIZE:
                print(f"  ⚠ Image too large: {content_length} bytes")
                return None
            
            # Download with size limit
            content = b''
            for chunk in response.iter_content(chunk_size=8192):
                content += chunk
                if len(content) > MAX_IMAGE_SIZE:
                    print(f"  ⚠ Image exceeds size limit")
                    return None
            
            return io.BytesIO(content)
        return None
    except Exception as e:
        print(f"Error downloading image: {e}")
        return None


def calculate_title_font_size(text, max_width_inches=8.5, bold=True):
    """
    Calculate optimal font size for title to fit in one line.
    Tries sizes from 40pt down to 24pt.
    Returns the largest font size that fits the text in one line.
    
    Approximate calculation: 1 character ≈ 0.6 * font_size_pt / 72 inches (for bold text)
    """
    # Font sizes to try in descending order
    font_sizes = [40, 36, 32, 28, 24]
    
    # Approximate character width factor for bold fonts (empirical)
    # For bold fonts: ~0.55-0.65 of font size in points
    char_width_factor = 0.6 if bold else 0.5
    
    for font_size in font_sizes:
        # Estimate text width in inches
        # char_width_in_points = font_size * char_width_factor
        # char_width_in_inches = char_width_in_points / 72
        estimated_width = len(text) * (font_size * char_width_factor / 72)
        
        if estimated_width <= max_width_inches:
            print(f"  📏 Title font size: {font_size}pt (estimated width: {estimated_width:.2f}in vs max {max_width_inches}in)")
            return font_size
    
    # If even 24pt doesn't fit, return 24pt anyway (minimum)
    print(f"  ⚠ Title very long, using minimum font size: 24pt")
    return 24


# Theme color configurations for presentations
PRESENTATION_THEMES = {
    'light': {
        'background': RGBColor(245, 245, 250),
        'title_slide_bg': RGBColor(15, 25, 45),
        'content_slide_bg': RGBColor(245, 245, 250),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(30, 60, 120),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(40, 40, 40),
        'accent_color': RGBColor(30, 60, 180)
    },
    'dark': {
        'background': RGBColor(30, 30, 30),
        'title_slide_bg': RGBColor(15, 15, 25),
        'content_slide_bg': RGBColor(30, 30, 30),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(187, 134, 252),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(224, 224, 224),
        'accent_color': RGBColor(3, 218, 198)
    },
    'modern': {
        'background': RGBColor(250, 250, 250),
        'title_slide_bg': RGBColor(30, 30, 50),
        'content_slide_bg': RGBColor(250, 250, 250),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(79, 70, 229),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(15, 23, 42),
        'accent_color': RGBColor(124, 58, 237)
    },
    'casual': {
        'background': RGBColor(255, 245, 247),
        'title_slide_bg': RGBColor(76, 69, 105),
        'content_slide_bg': RGBColor(255, 245, 247),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(255, 107, 157),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(51, 51, 51),
        'accent_color': RGBColor(255, 160, 122)
    },
    'classic': {
        'background': RGBColor(236, 240, 241),
        'title_slide_bg': RGBColor(28, 38, 50),
        'content_slide_bg': RGBColor(236, 240, 241),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(44, 62, 80),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(44, 62, 80),
        'accent_color': RGBColor(52, 73, 94)
    },
    'futuristic': {
        'background': RGBColor(10, 14, 39),
        'title_slide_bg': RGBColor(10, 14, 39),
        'content_slide_bg': RGBColor(10, 14, 39),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(0, 212, 255),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(255, 255, 255),
        'accent_color': RGBColor(255, 0, 255)
    },
    'gradient': {
        'background': RGBColor(254, 249, 255),
        'title_slide_bg': RGBColor(79, 30, 85),
        'content_slide_bg': RGBColor(254, 249, 255),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(240, 147, 251),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(51, 51, 51),
        'accent_color': RGBColor(79, 172, 254)
    },
    'glassmorphism': {
        'background': RGBColor(102, 126, 234),
        'title_slide_bg': RGBColor(26, 30, 74),
        'content_slide_bg': RGBColor(102, 126, 234),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(255, 255, 255),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(255, 255, 255),
        'accent_color': RGBColor(255, 255, 255)
    },
    'nature': {
        'background': RGBColor(241, 250, 238),
        'title_slide_bg': RGBColor(29, 67, 50),
        'content_slide_bg': RGBColor(241, 250, 238),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(45, 106, 79),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(27, 67, 50),
        'accent_color': RGBColor(82, 183, 136)
    },
    'vivid': {
        'background': RGBColor(255, 252, 242),
        'title_slide_bg': RGBColor(33, 5, 17),
        'content_slide_bg': RGBColor(255, 252, 242),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(255, 0, 110),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(33, 37, 41),
        'accent_color': RGBColor(251, 86, 7)
    },
    'business': {
        'background': RGBColor(248, 250, 252),
        'title_slide_bg': RGBColor(15, 25, 50),
        'content_slide_bg': RGBColor(248, 250, 252),
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(30, 58, 138),
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(15, 23, 42),
        'accent_color': RGBColor(14, 165, 233)
    },
    'heroic': {
        # HEROIC_MINIMAL: Clean minimalist design with strategic metaphors
        'background': RGBColor(255, 255, 255),  # Pure white
        'title_slide_bg': RGBColor(45, 55, 72),  # #2D3748 - Dark slate for title
        'content_slide_bg': RGBColor(255, 255, 255),  # White for content
        'title_color_first_last': RGBColor(255, 255, 255),  # White on dark
        'title_color_content': RGBColor(26, 32, 44),  # #1A202C - Almost black
        'content_color_first_last': RGBColor(255, 255, 255),  # White
        'content_color_content': RGBColor(26, 32, 44),  # #1A202C
        'accent_color': RGBColor(66, 153, 225),  # #4299E1 - Blue accent
        'icon_color': RGBColor(74, 85, 104),  # #4A5568 - Gray for icons
        'metaphor_overlay_color': RGBColor(45, 55, 72),  # #2D3748
        'style': 'heroic_minimal',  # Special flag
        'metaphor_percentage': 40  # 40% slides get metaphorical images
    },
    'minimal': {
        # Updated MINIMAL: Even cleaner with more air
        'background': RGBColor(248, 250, 252),  # #F8FAFC - Soft white
        'title_slide_bg': RGBColor(0, 0, 0),  # Pure black
        'content_slide_bg': RGBColor(248, 250, 252),  # Soft white
        'title_color_first_last': RGBColor(255, 255, 255),
        'title_color_content': RGBColor(26, 32, 44),  # #1A202C
        'content_color_first_last': RGBColor(255, 255, 255),
        'content_color_content': RGBColor(26, 32, 44),
        'accent_color': RGBColor(74, 85, 104),  # #4A5568 - Subtle gray
        'icon_color': RGBColor(74, 85, 104),
        'style': 'minimal_clean',
        'metaphor_percentage': 10  # Only 10% metaphorical images
    }
}

def filter_quiz_and_assessment_slides(slides_data):
    """
    Filter out quiz, self-assessment, and review question slides.
    These are not suitable for academic/professional presentations.
    
    Returns: (filtered_slides, removed_slides_info)
    """
    quiz_keywords = [
        'quiz', 'test', 'self-check', 'self-assessment', 'questions for review',
        'квиз', 'тест', 'самопроверка', 'вопросы для проверки', 'проверка знаний',
        'check your knowledge', 'knowledge check', 'проверьте себя',
        'review questions', 'повторение', 'practice questions'
    ]
    
    filtered = []
    removed = []
    
    for idx, slide in enumerate(slides_data):
        title = slide.get('title', '').lower()
        content = slide.get('content', '').lower()
        
        # Check if slide contains quiz/assessment keywords
        is_quiz_slide = any(kw in title or kw in content for kw in quiz_keywords)
        
        # Additional check: slides with many question marks (likely quiz)
        question_count = content.count('?')
        has_many_questions = question_count >= 3
        
        if is_quiz_slide or has_many_questions:
            removed.append({
                'index': idx,
                'title': slide.get('title', ''),
                'reason': 'quiz/self-assessment content' if is_quiz_slide else 'multiple questions detected'
            })
            print(f"  ❌ Removed slide {idx + 1}: '{slide.get('title', '')}' ({removed[-1]['reason']})")
        else:
            filtered.append(slide)
    
    return filtered, removed


def get_icon_unicode_for_slide(slide_title: str, slide_content: str) -> str:
    """
    Select appropriate line-style icon (Unicode) based on slide content.
    Returns Unicode character for thin line icons (Heroicons/Feather style).
    
    Icon Categories:
    - Idea/Goal: 💡 (lightbulb), 🧭 (compass)
    - Process: ⚙️ (gear), ➡️ (arrow)
    - Comparison: ⚖️ (scales), 📊 (chart)
    - Success: 🏆 (trophy), 📈 (growth)
    - Warning: ⚠️ (warning)
    - Information: ℹ️ (info), 📝 (document)
    - Target: 🎯 (target)
    - Time: ⏱️ (stopwatch), 📅 (calendar)
    - People: 👥 (users), 🤝 (handshake)
    - Tools: 🔧 (wrench), 🛠️ (tools)
    - Security: 🔒 (lock), 🛡️ (shield)
    - Communication: 💬 (speech), 📧 (email)
    """
    combined_text = (slide_title + " " + slide_content).lower()
    
    # Idea/Innovation/Goal/Vision
    if any(word in combined_text for word in ['idea', 'innovation', 'vision', 'идея', 'инновация', 'визия', 'creative', 'творческий']):
        return "💡"  # Lightbulb
    
    # Direction/Strategy/Goal/Compass
    if any(word in combined_text for word in ['direction', 'strategy', 'goal', 'compass', 'navigate', 'направление', 'стратегия', 'цель', 'навигация']):
        return "🧭"  # Compass
    
    # Target/Focus/Objective
    if any(word in combined_text for word in ['target', 'focus', 'objective', 'aim', 'цель', 'фокус', 'задача']):
        return "🎯"  # Target
    
    # Process/System/Mechanism/Work
    if any(word in combined_text for word in ['process', 'system', 'mechanism', 'workflow', 'operation', 'процесс', 'система', 'механизм', 'работа']):
        return "⚙️"  # Gear
    
    # Growth/Success/Achievement/Increase
    if any(word in combined_text for word in ['growth', 'increase', 'success', 'achievement', 'improve', 'рост', 'успех', 'достижение', 'увеличение']):
        return "📈"  # Growth chart
    
    # Award/Trophy/Win/Victory
    if any(word in combined_text for word in ['award', 'trophy', 'win', 'victory', 'champion', 'награда', 'победа', 'чемпион']):
        return "🏆"  # Trophy
    
    # Comparison/Analysis/Balance
    if any(word in combined_text for word in ['compare', 'comparison', 'balance', 'versus', 'analysis', 'сравнение', 'анализ', 'баланс']):
        return "⚖️"  # Scales
    
    # Data/Chart/Statistics/Metrics
    if any(word in combined_text for word in ['data', 'chart', 'statistics', 'metrics', 'analytics', 'данные', 'статистика', 'метрики']):
        return "📊"  # Bar chart
    
    # Warning/Risk/Alert/Danger
    if any(word in combined_text for word in ['warning', 'risk', 'alert', 'danger', 'caution', 'предупреждение', 'риск', 'опасность']):
        return "⚠️"  # Warning
    
    # Time/Schedule/Deadline
    if any(word in combined_text for word in ['time', 'schedule', 'deadline', 'timeline', 'время', 'график', 'срок']):
        return "⏱️"  # Stopwatch
    
    # Calendar/Date/Event/Plan
    if any(word in combined_text for word in ['calendar', 'date', 'event', 'plan', 'schedule', 'календарь', 'дата', 'событие', 'план']):
        return "📅"  # Calendar
    
    # Team/People/Collaboration/Users
    if any(word in combined_text for word in ['team', 'people', 'collaboration', 'users', 'group', 'команда', 'люди', 'сотрудничество']):
        return "👥"  # Users
    
    # Partnership/Agreement/Handshake
    if any(word in combined_text for word in ['partnership', 'agreement', 'cooperation', 'alliance', 'партнерство', 'соглашение', 'сотрудничество']):
        return "🤝"  # Handshake
    
    # Tools/Build/Development
    if any(word in combined_text for word in ['tool', 'build', 'development', 'construct', 'инструмент', 'создание', 'разработка']):
        return "🔧"  # Wrench
    
    # Security/Protection/Safe
    if any(word in combined_text for word in ['security', 'protection', 'safe', 'secure', 'protect', 'безопасность', 'защита']):
        return "🔒"  # Lock
    
    # Communication/Message/Discussion
    if any(word in combined_text for word in ['communication', 'message', 'discussion', 'talk', 'коммуникация', 'сообщение', 'обсуждение']):
        return "💬"  # Speech bubble
    
    # Document/File/Report
    if any(word in combined_text for word in ['document', 'file', 'report', 'paper', 'документ', 'файл', 'отчет']):
        return "📝"  # Document
    
    # Default: Info icon
    return "ℹ️"  # Info


def should_use_metaphorical_image(slide_index: int, total_slides: int, slide_title: str, slide_content: str, metaphor_percentage: int) -> tuple[bool, str | None]:
    """
    Determine if a slide should use metaphorical image instead of icon.
    Returns: (use_metaphor, metaphor_keyword)
    
    Metaphorical images for key moments:
    - Compass: direction, strategy, navigation
    - Phoenix: rebirth, transformation, renewal
    - Door: opportunity, opening, entrance, beginning
    - Road: journey, path, progress
    - Fire: passion, energy, transformation
    - Mountain: challenge, achievement, peak
    - Bridge: connection, transition, crossing
    - Lighthouse: guidance, vision, clarity
    - Sunrise: beginning, hope, new start
    - Keys: solution, access, unlock
    """
    combined_text = (slide_title + " " + slide_content).lower()
    
    # Calculate if this slide should get a metaphor based on percentage
    # Key slides (first, last, middle) have higher priority
    is_key_slide = (slide_index == 0 or slide_index == total_slides - 1 or slide_index == total_slides // 2)
    
    # Threshold calculation: key slides more likely to get metaphors
    if is_key_slide:
        should_get_metaphor = (slide_index % max(1, int(100 / (metaphor_percentage * 1.5)))) == 0
    else:
        should_get_metaphor = (slide_index % max(1, int(100 / metaphor_percentage))) == 0
    
    if not should_get_metaphor:
        return False, None
    
    # Check for metaphorical keywords
    metaphor_map = {
        'compass': ['direction', 'strategy', 'navigate', 'path', 'way', 'course', 'направление', 'стратегия', 'навигация', 'путь'],
        'phoenix fire': ['rebirth', 'transformation', 'renewal', 'rise', 'resurrect', 'возрождение', 'трансформация', 'обновление'],
        'open door opportunity': ['opportunity', 'opening', 'entrance', 'beginning', 'start', 'door', 'возможность', 'начало', 'вход'],
        'road journey path': ['journey', 'road', 'progress', 'ahead', 'forward', 'путешествие', 'дорога', 'прогресс', 'вперед'],
        'bonfire flames': ['passion', 'energy', 'fire', 'burn', 'ignite', 'страсть', 'энергия', 'огонь'],
        'mountain peak summit': ['challenge', 'achievement', 'peak', 'summit', 'climb', 'overcome', 'вызов', 'достижение', 'вершина'],
        'bridge connection': ['connection', 'bridge', 'link', 'connect', 'transition', 'связь', 'мост', 'переход'],
        'lighthouse guidance': ['guidance', 'vision', 'clarity', 'light', 'beacon', 'руководство', 'визия', 'ясность'],
        'sunrise dawn': ['beginning', 'hope', 'new', 'dawn', 'start', 'morning', 'начало', 'надежда', 'рассвет'],
        'golden key solution': ['solution', 'key', 'unlock', 'access', 'answer', 'решение', 'ключ', 'ответ']
    }
    
    # Find matching metaphor
    for metaphor_query, keywords in metaphor_map.items():
        if any(keyword in combined_text for keyword in keywords):
            return True, metaphor_query
    
    # No specific metaphor match, but slide was selected - use generic inspiring image
    return False, None


def create_presentation(topic, slides_data, theme='light', presentation_type='business', user_id=None):
    """
    Create PowerPoint presentation with text and images.
    Now includes:
    - Quiz/self-assessment slide filtering
    - Intelligent image search based on SLIDE CONTENT (not just topic)
    - Duplicate image prevention (within presentation + user history)
    - Dynamic font sizing to prevent text overflow
    
    Args:
        topic: Presentation topic
        slides_data: List of slide dicts with 'title' and 'content'
        theme: Visual theme (light/dark)
        presentation_type: Type (business/scientific/general)
        user_id: User ID for tracking image usage (optional)
    """
    print(f"\n{'#'*60}")
    print(f"# Creating presentation: {topic}")
    print(f"# Total slides (before filtering): {len(slides_data)}")
    print(f"# Theme: {theme}")
    print(f"# Type: {presentation_type}")
    if user_id:
        print(f"# User ID: {user_id}")
    print(f"{'#'*60}\n")
    
    # Filter out quiz/self-assessment slides
    print(f"\n📦 FILTERING QUIZ/ASSESSMENT SLIDES...")
    slides_data, removed_slides = filter_quiz_and_assessment_slides(slides_data)
    
    if removed_slides:
        print(f"\n⚠️ Removed {len(removed_slides)} quiz/assessment slide(s):")
        for r in removed_slides:
            print(f"  - Slide {r['index'] + 1}: '{r['title']}' ({r['reason']})")
    
    print(f"\n✅ Final slide count: {len(slides_data)} slides")
    print(f"{'='*60}\n")
    
    # Get theme configuration
    theme_config = PRESENTATION_THEMES.get(theme, PRESENTATION_THEMES['light'])
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # ============================================================================
    # DUPLICATE IMAGE PREVENTION SYSTEM
    # ============================================================================
    # Track used images in TWO ways:
    # 1. Within this presentation (used_images set)
    # 2. Across user's history (exclude_images from database)
    
    used_images = set()  # Images used in this presentation
    exclude_images = []  # Images to exclude (from user history)
    
    if user_id:
        # Get user's recently used images to avoid duplicates
        exclude_images = get_used_images_for_user(user_id, limit=100)
        if exclude_images:
            print(f"📊 Loaded {len(exclude_images)} previously used images for user {user_id}")
            print(f"   → Will avoid these in image search to prevent duplicates\n")
    
    for idx, slide_data in enumerate(slides_data):
        # Add a blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background color based on theme
        background = slide.background
        fill = background.fill
        fill.solid()
        
        is_title_slide = (idx == 0)
        is_last_slide = (idx == len(slides_data) - 1)
        
        if is_title_slide or is_last_slide:
            fill.fore_color.rgb = theme_config['title_slide_bg']
        else:
            fill.fore_color.rgb = theme_config['content_slide_bg']
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(8.5), Inches(0.8)
        )
        # Title style
        title_frame = title_box.text_frame
        title_frame.word_wrap = False  # NO line breaks - single line only
        title_frame.text = slide_data['title']
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.name = 'Roboto'
        
        # Calculate optimal font size to fit title in one line
        optimal_font_size = calculate_title_font_size(
            text=slide_data['title'],
            max_width_inches=8.5,
            bold=True
        )
        
        if is_title_slide or is_last_slide:
            title_para.font.size = Pt(optimal_font_size)
            title_para.font.bold = True
            title_para.font.color.rgb = theme_config['title_color_first_last']
        else:
            title_para.font.size = Pt(optimal_font_size)
            title_para.font.bold = True
            title_para.font.color.rgb = theme_config['title_color_content']
        
        # Add accent element for content slides based on theme
        if not (is_title_slide or is_last_slide):
            try:
                slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    Inches(0.3), Inches(0.3), Inches(0.1), Inches(5.0)
                ).fill.fore_color.rgb = theme_config['accent_color']
            except Exception as e:
                print(f"  ⚠ Failed to add left bar: {e}")
        
        # ============================================================================
        # SMART IMAGE SEARCH PER SLIDE
        # ============================================================================
        # Uses slide-specific content analysis for better image matching
        # Prevents duplicates within presentation AND across user history
        
        print(f"\n[Slide {idx + 1}/{len(slides_data)}] {slide_data['title']}")
        print(f"  Content: {slide_data['content'][:60]}...")
        
        # Combine within-presentation used images + user history
        all_exclude_images = list(used_images) + exclude_images
        
        # Extract search_keyword and image_prompt from slide_data (supports both modes)
        search_keyword = slide_data.get('search_keyword', None)
        image_prompt = slide_data.get('image_prompt', None)
        
        # Route based on USE_IMAGE_PROMPT flag
        if not USE_IMAGE_PROMPT:
            # LEGACY MODE: Use search_keyword
            image_data, image_url, query_used = search_image_legacy_mode(
                slide_title=slide_data['title'],
                slide_content=slide_data.get('content', ''),
                main_topic=topic,
                exclude_images=all_exclude_images,
                presentation_type=presentation_type,
                search_keyword=search_keyword,  # LLM-generated in English
                language=None  # Auto-detect
            )
        else:
            # ADVANCED MODE: Use image_prompt
            image_data, image_url, query_used = search_image_advanced_mode(
                slide_title=slide_data['title'],
                slide_content=slide_data.get('content', ''),
                main_topic=topic,
                exclude_images=all_exclude_images,
                presentation_type=presentation_type,
                image_prompt=image_prompt,  # LLM-generated description
                language=None  # Auto-detect
            )
        
        if image_data and image_url:
            # Mark image as used in this presentation
            used_images.add(image_url)
            
            # Track in database for future duplicate prevention
            if user_id:
                add_used_image(user_id, image_url, query_used or slide_data['title'])
            
            try:
                # Add image on the right side
                slide.shapes.add_picture(
                    image_data,
                    Inches(5.5), Inches(1.3),
                    width=Inches(4),
                    height=Inches(3.5)
                )
                print(f"  ✅ Image added successfully (unique, query: '{query_used}')")
            except Exception as e:
                print(f"  ✗ Error adding image to slide: {e}")
        else:
            print(f"  ⚠️ Continuing without image (no suitable unique image found)")
        
        # Add content text with improved overflow handling
        content_text = slide_data['content']
        
        # For first/last slides: more aggressive length limiting to prevent overflow
        if is_title_slide or is_last_slide:
            max_chars = 350  # Shorter limit for title/conclusion slides
            if len(content_text) > max_chars:
                content_text = content_text[:max_chars] + "..."
                print(f"  ✂️ Content trimmed: {len(slide_data['content'])} → {len(content_text)} chars (title/last slide)")
        else:
            # Content slides can have more text
            if len(content_text) > 500:
                content_text = content_text[:500] + "..."
                print(f"  ✂️ Content trimmed: {len(slide_data['content'])} → {len(content_text)} chars")
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4),
            Inches(4.8), Inches(3.6)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content_text
        
        # Format content text based on theme and slide position
        # Dynamic font size based on content length and slide type
        content_length = len(content_text)
        
        # First/last slides need smaller fonts to prevent overflow
        if is_title_slide or is_last_slide:
            if content_length > 280:
                base_font_size = 16
            elif content_length > 220:
                base_font_size = 17
            elif content_length > 160:
                base_font_size = 18
            else:
                base_font_size = 20
            print(f"  🔤 Title/Last slide font: {base_font_size}pt (length: {content_length})")
        else:
            # Content slides
            if content_length > 250:
                base_font_size = 14
            elif content_length > 180:
                base_font_size = 15
            else:
                base_font_size = 16
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.name = 'Roboto'
            paragraph.font.size = Pt(base_font_size if not (is_title_slide or is_last_slide) else 20)
            if is_title_slide or is_last_slide:
                paragraph.font.color.rgb = theme_config['content_color_first_last']
            else:
                paragraph.font.color.rgb = theme_config['content_color_content']
            paragraph.space_after = Pt(10)
            paragraph.line_spacing = 1.2
        
        print(f"\n{'='*60}")
        print(f"✓ Slide {idx + 1} created successfully")
        print(f"  Title: {slide_data['title']}")
        print(f"  Content length: {len(slide_data['content'])} characters")
        print(f"{'='*60}")
    
    # Save presentation
    filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    
    # Cleanup old used images for this user (keep last 100)
    if user_id:
        cleanup_old_used_images(user_id, keep_count=100)
    
    print(f"\n{'#'*60}")
    print(f"# ✓ Presentation created successfully!")
    print(f"# File: {filename}")
    print(f"# Location: {filepath}")
    print(f"# Unique images used: {len(used_images)}")
    if user_id:
        print(f"# Total images tracked for user: {len(exclude_images) + len(used_images)}")
    print(f"{'#'*60}\n")
    
    return filepath


@app.route('/')
def index():
    """
    Render main page
    """
    return render_template('index.html')


@app.route('/pricing')
def pricing():
    """
    Render pricing page with Stripe payment options
    """
    return render_template('pricing.html')


# ============================================================================
# STRIPE PAYMENT ROUTES
# ============================================================================

@app.route('/api/create-checkout-session', methods=['POST'])
@login_required
def create_checkout_session():
    """
    Create Stripe Checkout Session for payment
    Accepts: plan_type (one_time, subscription, pro, premium)
    Returns: sessionId and checkout URL
    """
    try:
        # Check if Stripe is configured
        if not STRIPE_SECRET_KEY:
            return jsonify({
                'error': 'Payment system not configured',
                'message': 'Stripe is not configured on the server'
            }), 500
        
        data = request.json
        plan_type = data.get('plan_type', 'one_time')
        
        # Get user email for Stripe customer
        user_email = current_user.email if hasattr(current_user, 'email') else None
        
        if not user_email:
            return jsonify({'error': 'User email not found'}), 400
        
        # Define pricing based on plan type
        # NOTE: You need to create these prices in your Stripe Dashboard
        # and replace with actual price IDs
        price_configs = {
            'one_time': {
                'mode': 'payment',
                'price_data': {
                    'currency': 'usd',
                    'unit_amount': 999,  # $9.99
                    'product_data': {
                        'name': 'AI SlideRush - Single Purchase',
                        'description': 'One-time access to create presentations',
                    },
                },
                'quantity': 1,
                'plan_name': 'one_time'
            },
            'subscription': {
                'mode': 'subscription',
                'price_data': {
                    'currency': 'usd',
                    'unit_amount': 1999,  # $19.99/month
                    'recurring': {'interval': 'month'},
                    'product_data': {
                        'name': 'AI SlideRush - Monthly Subscription',
                        'description': 'Unlimited presentations per month',
                    },
                },
                'quantity': 1,
                'plan_name': 'subscription'
            },
            'pro': {
                'mode': 'payment',
                'price_data': {
                    'currency': 'usd',
                    'unit_amount': 1999,  # $19.99
                    'product_data': {
                        'name': 'AI SlideRush - Pro Plan',
                        'description': 'Pro features with advanced customization',
                    },
                },
                'quantity': 1,
                'plan_name': 'pro'
            },
            'premium': {
                'mode': 'subscription',
                'price_data': {
                    'currency': 'usd',
                    'unit_amount': 4999,  # $49.99/month
                    'recurring': {'interval': 'month'},
                    'product_data': {
                        'name': 'AI SlideRush - Premium Subscription',
                        'description': 'Unlimited presentations with priority support',
                    },
                },
                'quantity': 1,
                'plan_name': 'premium'
            }
        }
        
        if plan_type not in price_configs:
            return jsonify({'error': f'Invalid plan type: {plan_type}'}), 400
        
        config = price_configs[plan_type]
        
        # Create Stripe Checkout Session
        try:
            checkout_session = stripe.checkout.Session.create(
                payment_method_types=['card'],
                line_items=[{
                    'price_data': config['price_data'],
                    'quantity': config['quantity'],
                }],
                mode=config['mode'],
                success_url=request.host_url + 'dashboard?payment=success',
                cancel_url=request.host_url + 'dashboard?payment=cancelled',
                customer_email=user_email,
                client_reference_id=str(current_user.id),  # Store user ID for webhook
                metadata={
                    'user_id': str(current_user.id),
                    'user_email': user_email,
                    'plan_type': config['plan_name']
                }
            )
            
            print(f"✅ Stripe Checkout Session created: {checkout_session.id}")
            print(f"   → User: {user_email} (ID: {current_user.id})")
            print(f"   → Plan: {plan_type}")
            print(f"   → Amount: ${config['price_data']['unit_amount'] / 100:.2f}")
            
            return jsonify({
                'success': True,
                'sessionId': checkout_session.id,
                'url': checkout_session.url,
                'plan_type': plan_type
            })
            
        except stripe.error.StripeError as e:
            print(f"❌ Stripe error: {e}")
            return jsonify({
                'error': 'Payment session creation failed',
                'message': str(e)
            }), 500
            
    except Exception as e:
        print(f"❌ Error creating checkout session: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/stripe/webhook', methods=['POST'])
def stripe_webhook():
    """
    Stripe webhook handler for payment events
    Verifies webhook signature and processes checkout.session.completed
    """
    payload = request.data
    sig_header = request.headers.get('Stripe-Signature')
    
    # Verify webhook signature (if secret is configured)
    if STRIPE_WEBHOOK_SECRET:
        try:
            event = stripe.Webhook.construct_event(
                payload, sig_header, STRIPE_WEBHOOK_SECRET
            )
            print(f"✅ Webhook signature verified: {event['type']}")
        except ValueError as e:
            # Invalid payload
            print(f"❌ Webhook error: Invalid payload - {e}")
            return jsonify({'error': 'Invalid payload'}), 400
        except stripe.error.SignatureVerificationError as e:
            # Invalid signature
            print(f"❌ Webhook error: Invalid signature - {e}")
            return jsonify({'error': 'Invalid signature'}), 400
    else:
        # No signature verification (INSECURE - only for development)
        print("⚠️ Webhook signature verification SKIPPED (no STRIPE_WEBHOOK_SECRET)")
        try:
            event = json.loads(payload)
        except json.JSONDecodeError as e:
            print(f"❌ Webhook error: Invalid JSON - {e}")
            return jsonify({'error': 'Invalid JSON'}), 400
    
    # Handle the event
    event_type = event['type']
    print(f"\n{'='*60}")
    print(f"📥 Stripe Webhook Event: {event_type}")
    print(f"{'='*60}")
    
    if event_type == 'checkout.session.completed':
        session = event['data']['object']
        
        # Extract customer information
        customer_email = session.get('customer_details', {}).get('email')
        client_reference_id = session.get('client_reference_id')  # User ID
        metadata = session.get('metadata', {})
        plan_type = metadata.get('plan_type', 'one_time')
        stripe_customer_id = session.get('customer')
        
        print(f"Payment completed:")
        print(f"  → Email: {customer_email}")
        print(f"  → User ID: {client_reference_id}")
        print(f"  → Plan: {plan_type}")
        print(f"  → Customer ID: {stripe_customer_id}")
        
        # Find user in database
        user_id = None
        if client_reference_id:
            user_id = client_reference_id
        elif customer_email:
            # Fallback: find user by email
            try:
                conn = sqlite3.connect(DB_PATH)
                conn.row_factory = sqlite3.Row
                cursor = conn.cursor()
                cursor.execute('SELECT id FROM users WHERE email = ?', (customer_email,))
                user_row = cursor.fetchone()
                conn.close()
                if user_row:
                    user_id = user_row['id']
                    print(f"  → Found user by email: ID {user_id}")
            except Exception as e:
                print(f"❌ Error finding user by email: {e}")
        
        if not user_id:
            print(f"❌ Cannot find user for payment (email: {customer_email})")
            return jsonify({'error': 'User not found'}), 400
        
        # Update user in database
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            
            # Determine subscription status based on plan type
            is_subscription = plan_type in ['subscription', 'premium']
            subscription_status = 'active' if is_subscription else 'active'
            
            cursor.execute('''
                UPDATE users 
                SET status = 'active',
                    stripe_customer_id = ?,
                    subscription_plan = ?,
                    subscription_status = ?
                WHERE id = ?
            ''', (stripe_customer_id, plan_type, subscription_status, user_id))
            
            conn.commit()
            affected_rows = cursor.rowcount
            conn.close()
            
            if affected_rows > 0:
                print(f"✅ User {user_id} updated successfully:")
                print(f"   → Status: active")
                print(f"   → Plan: {plan_type}")
                print(f"   → Subscription status: {subscription_status}")
                print(f"   → Stripe customer: {stripe_customer_id}")
            else:
                print(f"⚠️ No user updated (user_id={user_id} not found)")
                
        except Exception as e:
            print(f"❌ Error updating user after payment: {e}")
            return jsonify({'error': 'Database update failed'}), 500
    
    elif event_type == 'customer.subscription.deleted':
        # Handle subscription cancellation
        subscription = event['data']['object']
        stripe_customer_id = subscription.get('customer')
        
        print(f"Subscription cancelled for customer: {stripe_customer_id}")
        
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            cursor.execute('''
                UPDATE users 
                SET subscription_status = 'cancelled',
                    subscription_plan = 'free'
                WHERE stripe_customer_id = ?
            ''', (stripe_customer_id,))
            conn.commit()
            conn.close()
            print(f"✅ User subscription status updated to cancelled")
        except Exception as e:
            print(f"❌ Error updating subscription status: {e}")
    
    elif event_type == 'invoice.payment_failed':
        # Handle failed payment
        invoice = event['data']['object']
        stripe_customer_id = invoice.get('customer')
        
        print(f"⚠️ Payment failed for customer: {stripe_customer_id}")
        
        # Optionally update user status
        # (you might want to give them a grace period before blocking)
    
    else:
        print(f"ℹ️ Unhandled event type: {event_type}")
    
    print(f"{'='*60}\n")
    
    return jsonify({'success': True}), 200


@app.route('/api/create-presentation', methods=['POST'])
def create_presentation_api():
    """
    API endpoint to create presentation
    Now includes payment status verification
    """
    try:
        # ====================================================================
        # 🧠 MANDATORY CLIP STATUS CHECK BEFORE GENERATION
        # ====================================================================
        print("\n" + "="*70)
        print("🧠 CLIP STATUS PRE-FLIGHT CHECK")
        print("="*70)
        
        if not CLIP_AVAILABLE:
            print("❌ CRITICAL: CLIP not available - cannot generate presentation")
            print("   This should never happen if server started correctly!")
            print("="*70 + "\n")
            return jsonify({
                'error': 'CLIP service unavailable',
                'message': 'Image matching service is not available. Please contact administrator.',
                'clip_status': 'unavailable'
            }), 500
        
        # Verify CLIP model is actually loaded
        from services import clip_client
        if clip_client._clip_model is None:
            print("❌ CRITICAL: CLIP model is None - system in invalid state")
            print("="*70 + "\n")
            return jsonify({
                'error': 'CLIP model not loaded',
                'message': 'Image matching service failed to initialize. Please restart server.',
                'clip_status': 'not_loaded'
            }), 500
        
        print("✅ CLIP Status: READY")
        print(f"   → Model: {clip_client._clip_model.__class__.__name__}")
        print(f"   → Device: {clip_client._device}")
        print(f"   → Cache size: {len(clip_client._image_embedding_cache)}")
        print("="*70 + "\n")
        
        # ====================================================================
        # Continue with normal request processing
        # ====================================================================
        data = request.json
        topic = data.get('topic', '').strip()
        num_slides = data.get('num_slides', 5)
        language = data.get('language', 'en')  # Get language from frontend
        theme = data.get('theme', 'light')  # Get theme from frontend
        presentation_type = data.get('presentation_type', 'business')  # Get presentation type
        
        # Validation
        if not topic:
            return jsonify({'error': 'Topic is required'}), 400
        
        # ============================================================================
        # FREE CREDITS & PAYMENT VERIFICATION
        # ============================================================================
        # Order of checks:
        # 1. Check if payments are enabled (dev mode bypass)
        # 2. Block if user is 'blocked'
        # 3. Allow if user has free_credits > 0 (no Stripe check needed)
        # 4. Only check Stripe if free_credits == 0
        
        # DEV MODE: Skip all payment checks if PAYMENTS_ENABLED=false
        if not PAYMENTS_ENABLED:
            print(f"[DEV] ⚠️ Payments disabled, skipping payment check for user {current_user.id if current_user.is_authenticated else 'anonymous'}")
            request.using_free_credit = False  # Not using credit in dev mode
        elif current_user.is_authenticated:
            try:
                conn = sqlite3.connect(DB_PATH)
                conn.row_factory = sqlite3.Row
                cursor = conn.cursor()
                cursor.execute(
                    'SELECT status, subscription_plan, subscription_status, free_credits FROM users WHERE id = ?',
                    (current_user.id,)
                )
                user_row = cursor.fetchone()
                conn.close()
                
                if user_row:
                    # ========================================
                    # STEP 1: Check if user is blocked
                    # ========================================
                    if user_row['status'] == 'blocked':
                        print(f"⛔ User {current_user.id} is BLOCKED - cannot create presentations")
                        return jsonify({
                            'error': 'Account blocked',
                            'message': 'Your account has been blocked. Please contact support.',
                            'requires_payment': False
                        }), 403
                    
                    # ========================================
                    # STEP 2: Handle free_credits (backward compatibility)
                    # ========================================
                    free_credits = user_row['free_credits']
                    
                    # Backward compatibility: if free_credits is NULL for existing users, initialize to 3
                    if free_credits is None:
                        print(f"🔄 User {current_user.id} has NULL free_credits - initializing to 3")
                        try:
                            conn = sqlite3.connect(DB_PATH)
                            cursor = conn.cursor()
                            cursor.execute('UPDATE users SET free_credits = 3 WHERE id = ?', (current_user.id,))
                            conn.commit()
                            conn.close()
                            free_credits = 3
                            print(f"   → Initialized free_credits = 3 for user {current_user.id}")
                        except Exception as e:
                            print(f"⚠️ Error initializing free_credits: {e}")
                            free_credits = 0  # Fallback to 0 if update fails
                    
                    # ========================================
                    # STEP 3: Check free credits first (bypass Stripe)
                    # ========================================
                    if free_credits > 0:
                        print(f"🎁 User {current_user.id} using FREE CREDIT ({free_credits} remaining)")
                        print(f"   → Bypassing Stripe payment verification")
                        # Will decrement free_credits after successful presentation creation
                        # Store in session/variable to decrement later
                        request.using_free_credit = True
                    
                    # ========================================
                    # STEP 4: Only check Stripe if NO free credits
                    # ========================================
                    else:
                        print(f"💳 User {current_user.id} has 0 free credits - checking Stripe subscription")
                        subscription_status = user_row['subscription_status'] or 'inactive'
                        subscription_plan = user_row['subscription_plan'] or 'free'
                        
                        # Allow if: subscription_status='active' OR subscription_plan != 'free'
                        has_valid_subscription = (subscription_status == 'active') or (subscription_plan != 'free')
                        
                        if not has_valid_subscription:
                            print(f"⛔ User {current_user.id} requires payment:")
                            print(f"   → Plan: {subscription_plan}")
                            print(f"   → Status: {subscription_status}")
                            print(f"   → Free credits: {free_credits}")
                            return jsonify({
                                'error': 'Payment required',
                                'message': 'You have used all 3 free presentations. Please upgrade your plan to continue.',
                                'requires_payment': True,
                                'current_plan': subscription_plan,
                                'free_credits_remaining': 0
                            }), 403
                        
                        print(f"✅ User {current_user.id} has valid subscription - plan: {subscription_plan}, status: {subscription_status}")
                        request.using_free_credit = False
                        
            except Exception as e:
                print(f"⚠️ Error checking user payment status: {e}")
                # Continue anyway for backward compatibility
                request.using_free_credit = False
        
        # Normalize slides count to 5-10 range (enforced for 3 types)
        try:
            num_slides = int(num_slides)
        except (ValueError, TypeError):
            num_slides = 7  # Default to middle of range
        if num_slides < 5 or num_slides > 10:
            num_slides = max(5, min(10, num_slides))  # Clamp to 5-10
        
        # Check API keys
        if not OPENAI_API_KEY:
            return jsonify({'error': 'OpenAI API key not configured'}), 500
        
        if not PEXELS_API_KEY:
            return jsonify({'error': 'Pexels API key not configured'}), 500
        
        # Generate slide content in the selected language
        print(f"Generating content for topic: {topic}, slides: {num_slides}, language: {language}, type: {presentation_type}")
        slides_data = generate_slide_content_in_language(topic, num_slides, language, presentation_type)
        
        if not slides_data:
            # Use fallback slides in the selected language
            print("Using fallback slides in selected language")
            slides_data = create_fallback_slides(topic, num_slides, language)
            if not slides_data:
                return jsonify({'error': 'Failed to generate slide content'}), 502
        
        # Ensure we have the right number of slides
        slides_data = slides_data[:num_slides]
        
        # Create presentation with the selected theme and presentation type
        # Pass user_id for image duplicate tracking
        print("Creating presentation with theme:", theme, "type:", presentation_type)
        user_id_for_images = current_user.id if current_user.is_authenticated else None
        filepath = create_presentation(topic, slides_data, theme, presentation_type, user_id=user_id_for_images)
        filename = os.path.basename(filepath)
        
        # Save presentation to database if user is authenticated (already verified above)
        if current_user.is_authenticated:
            try:
                conn = sqlite3.connect(DB_PATH)
                cursor = conn.cursor()
                cursor.execute(
                    '''INSERT INTO presentations (user_id, topic, num_slides, filename, presentation_type) 
                       VALUES (?, ?, ?, ?, ?)''',
                    (current_user.id, topic, num_slides, filename, presentation_type)
                )
                conn.commit()
                print(f"✅ Presentation saved to database for user {current_user.id}")
                
                # ========================================
                # DECREMENT FREE CREDITS if used
                # ========================================
                if hasattr(request, 'using_free_credit') and request.using_free_credit:
                    cursor.execute(
                        'UPDATE users SET free_credits = free_credits - 1 WHERE id = ?',
                        (current_user.id,)
                    )
                    conn.commit()
                    
                    # Get updated credits count
                    cursor.execute('SELECT free_credits FROM users WHERE id = ?', (current_user.id,))
                    updated_row = cursor.fetchone()
                    credits_remaining = updated_row[0] if updated_row else 0
                    
                    print(f"🎁 FREE CREDIT USED - User {current_user.id} now has {credits_remaining} free presentations remaining")
                    if credits_remaining == 0:
                        print(f"   ⚠️ User has exhausted free credits - next presentation will require payment")
                
                conn.close()
            except Exception as e:
                print(f"⚠️ Error saving presentation to database: {e}")
        
        return jsonify({
            'success': True,
            'filename': filename,
            'slides': slides_data,
            'presentation_type': presentation_type
        })
        
    except Exception as e:
        print(f"Error in create_presentation_api: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/download/<filename>')
def download_presentation(filename):
    """
    Download generated presentation
    Security: Prevent path traversal attacks
    """
    try:
        # Security: Normalize path and prevent directory traversal
        filename = os.path.basename(filename)  # Remove any path components
        if '..' in filename or '/' in filename or '\\' in filename:
            return jsonify({'error': 'Invalid filename'}), 400
        
        filepath = os.path.join(OUTPUT_DIR, filename)
        
        # Security: Ensure file is within OUTPUT_DIR
        if not os.path.abspath(filepath).startswith(os.path.abspath(OUTPUT_DIR)):
            return jsonify({'error': 'Access denied'}), 403
        
        if not os.path.exists(filepath):
            return jsonify({'error': 'File not found'}), 404
        
        return send_file(
            filepath,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/presentation-types', methods=['GET'])
def get_presentation_types():
    """
    API endpoint to get all available presentation types
    """
    return jsonify({
        'success': True,
        'types': PRESENTATION_TYPES
    })


@app.route('/health', methods=['GET'])
@app.route('/api/health', methods=['GET'])
def health_check():
    """
    Health check endpoint for Railway deployment monitoring.
    Returns service status and configuration info.
    """
    import sys
    
    health_data = {
        'status': 'ok',
        'timestamp': datetime.now().isoformat(),
        'environment': 'railway' if IS_RAILWAY else 'local',
        'python_version': f"{sys.version_info.major}.{sys.version_info.minor}.{sys.version_info.micro}",
        'services': {
            'clip': 'enabled' if CLIP_AVAILABLE else 'disabled',
            'openai': 'configured' if OPENAI_API_KEY else 'missing',
            'pexels': 'configured' if PEXELS_API_KEY else 'missing',
            'firebase': 'configured' if os.getenv('FIREBASE_CREDENTIALS') else 'missing',
            'stripe': 'configured' if STRIPE_SECRET_KEY else 'missing'
        },
        'features': {
            'image_search': 'clip+keyword' if CLIP_AVAILABLE else 'keyword-only',
            'translation': TRANSLATION_ENABLED,
            'payments': PAYMENTS_ENABLED
        }
    }
    
    return jsonify(health_data), 200


@app.route('/api/test-clip', methods=['GET'])
def test_clip():
    """
    🧪 EMERGENCY DIAGNOSTIC: Test endpoint for CLIP performance.
    Tests image search with CLIP matching for a given text query.
    
    TARGET: <5 seconds response time for single query
    GOAL: <60 seconds for full presentation (10-15 slides)
    
    Query params:
        text: Search text (default: "banana plantain")
    
    Returns:
        JSON with comprehensive timing and diagnostic information
    """
    
    # Get query parameter
    text_query = request.args.get('text', 'banana plantain')
    
    print(f"\n{'='*70}")
    print(f"🧪 CLIP PERFORMANCE TEST - EMERGENCY DIAGNOSTIC")
    print(f"{'='*70}")
    print(f"🔍 Query: '{text_query}'")
    print(f"🎯 Target: <5 seconds for single test")
    
    overall_start = time.perf_counter()
    
    try:
        # DIAGNOSTIC 1: Check CLIP status
        print(f"\n🔧 STEP 1: CLIP Status Check")
        step1_start = time.perf_counter()
        
        print(f"   → CLIP_ENABLED: {CLIP_ENABLED}")
        print(f"   → CLIP_AVAILABLE: {CLIP_AVAILABLE}")
        print(f"   → CLIP_IMPORT_SUCCESS: {CLIP_IMPORT_SUCCESS}")
        
        # Import clip_client for detailed checks
        from services import clip_client
        
        if not CLIP_AVAILABLE:
            error_msg = "CLIP NOT AVAILABLE - Server should have failed to start!"
            print(f"\n❌ {error_msg}")
            print(f"   This is a CRITICAL ERROR - server should not be running!")
            return jsonify({
                'success': False,
                'error': error_msg,
                'elapsed_ms': (time.perf_counter() - overall_start) * 1000,
                'diagnostics': {
                    'clip_enabled': CLIP_ENABLED,
                    'clip_available': CLIP_AVAILABLE,
                    'clip_import_success': CLIP_IMPORT_SUCCESS,
                    'clip_model_loaded': clip_client._clip_model is not None
                }
            }), 500
        
        # Verify model is actually loaded
        if clip_client._clip_model is None:
            error_msg = "CLIP model is None - system in invalid state"
            print(f"\n❌ {error_msg}")
            return jsonify({
                'success': False,
                'error': error_msg,
                'elapsed_ms': (time.perf_counter() - overall_start) * 1000,
                'diagnostics': {
                    'clip_available': CLIP_AVAILABLE,
                    'clip_model_loaded': False
                }
            }), 500
        
        step1_time = (time.perf_counter() - step1_start) * 1000
        print(f"   ✅ CLIP is available and ready")
        print(f"   ⏱️  Time: {step1_time:.1f}ms")
        
        # DIAGNOSTIC 2: Fetch image candidates
        print(f"\n📸 STEP 2: Fetching Image Candidates (max 6)")
        step2_start = time.perf_counter()
        
        candidates = get_images(text_query, count=6)
        
        step2_time = (time.perf_counter() - step2_start) * 1000
        print(f"   → Found: {len(candidates) if candidates else 0} images")
        print(f"   ⏱️  Time: {step2_time:.1f}ms")
        
        if not candidates:
            error_msg = f"No image candidates found for query: '{text_query}'"
            print(f"\n❌ {error_msg}")
            return jsonify({
                'success': False,
                'error': error_msg,
                'elapsed_ms': (time.perf_counter() - overall_start) * 1000,
                'timing': {
                    'clip_check_ms': step1_time,
                    'fetch_candidates_ms': step2_time
                }
            }), 404
        
        # DIAGNOSTIC 3: Prepare candidates for CLIP
        print(f"\n🛠️  STEP 3: Preparing Candidates")
        step3_start = time.perf_counter()
        
        for i, candidate in enumerate(candidates):
            if 'description' not in candidate:
                candidate['description'] = (
                    candidate.get('attribution', '') or 
                    candidate.get('author', '') or 
                    text_query
                )
            print(f"   [{i+1}] {candidate.get('description', 'No description')[:40]}")
        
        step3_time = (time.perf_counter() - step3_start) * 1000
        print(f"   ⏱️  Time: {step3_time:.1f}ms")
        
        # DIAGNOSTIC 4: Run CLIP matching (THE CRITICAL STEP)
        print(f"\n🤖 STEP 4: CLIP SEMANTIC MATCHING (CRITICAL)")
        step4_start = time.perf_counter()
        
        best_image = clip_pick_best_image(
            slide_title=text_query,
            slide_content=f"Testing CLIP performance for query: {text_query}",
            image_candidates=candidates,
            exclude_images=[],
            similarity_threshold=0.0  # Soft mode for testing
        )
        
        step4_time = (time.perf_counter() - step4_start) * 1000
        print(f"\n   ⏱️  CLIP Time: {step4_time:.1f}ms ({step4_time/1000:.2f}s)")
        
        # Calculate total elapsed time
        total_elapsed = time.perf_counter() - overall_start
        total_ms = total_elapsed * 1000
        
        # Determine performance status
        if total_elapsed < 5.0:
            status = '✅ EXCELLENT'
            performance_level = 'excellent'
            status_emoji = '🎉'
        elif total_elapsed < 10.0:
            status = '✓ GOOD'
            performance_level = 'good'
            status_emoji = '👍'
        else:
            status = '⚠️ SLOW'
            performance_level = 'slow'
            status_emoji = '🐢'
        
        print(f"\n{'='*70}")
        print(f"{status_emoji} PERFORMANCE RESULTS: {status}")
        print(f"{'='*70}")
        print(f"⏱️  Total time: {total_ms:.1f}ms ({total_elapsed:.2f}s)")
        print(f"🎯 Target: <5000ms (5s)")
        print(f"{'✅ MEETS TARGET' if total_elapsed < 5.0 else '❌ EXCEEDS TARGET'}")
        print(f"\n📈 Breakdown:")
        print(f"   1. CLIP status check:    {step1_time:7.1f}ms")
        print(f"   2. Fetch candidates:     {step2_time:7.1f}ms")
        print(f"   3. Prepare data:         {step3_time:7.1f}ms")
        print(f"   4. CLIP matching:        {step4_time:7.1f}ms  ⭐ CRITICAL")
        print(f"   {'─'*40}")
        print(f"   TOTAL:                   {total_ms:7.1f}ms")
        print(f"\n💡 Estimated full presentation (15 slides):")
        estimated_full = (total_elapsed * 15)
        print(f"   {estimated_full:.1f}s ({estimated_full/60:.1f} min)")
        print(f"   {'✅ Under 60s target' if estimated_full < 60 else '❌ Over 60s target'}")
        print(f"{'='*70}\n")
        
        # Return detailed results
        result = {
            'success': True,
            'query': text_query,
            'performance': performance_level,
            'status': status,
            'elapsed_ms': round(total_ms, 1),
            'elapsed_sec': round(total_elapsed, 2),
            'target_sec': 5.0,
            'meets_target': total_elapsed < 5.0,
            'estimated_full_presentation': {
                'slides': 15,
                'estimated_sec': round(estimated_full, 1),
                'estimated_min': round(estimated_full / 60, 2),
                'meets_60s_target': estimated_full < 60
            },
            'timing': {
                'step1_clip_check_ms': round(step1_time, 1),
                'step2_fetch_candidates_ms': round(step2_time, 1),
                'step3_prepare_data_ms': round(step3_time, 1),
                'step4_clip_matching_ms': round(step4_time, 1),
                'total_ms': round(total_ms, 1)
            },
            'candidates': {
                'count': len(candidates),
                'max_allowed': 6,
                'optimized': True
            },
            'diagnostics': {
                'clip_enabled': CLIP_ENABLED,
                'clip_available': CLIP_AVAILABLE,
                'clip_import_success': CLIP_IMPORT_SUCCESS,
                'use_strict_filter': USE_STRICT_CLIP_FILTER,
                'similarity_threshold': CLIP_SIMILARITY_THRESHOLD
            }
        }
        
        if best_image:
            result['best_match'] = {
                'url': best_image.get('url', '')[:100],
                'similarity': best_image.get('_clip_similarity', 'N/A'),
                'source': best_image.get('source', 'Unknown'),
                'description': best_image.get('description', '')[:60]
            }
            print(f"✅ Best match found: {result['best_match']['description']}")
            print(f"   Similarity: {result['best_match']['similarity']}")
        else:
            print("⚠️ No best match returned (check CLIP logs above)")
        
        return jsonify(result)
        
    except Exception as e:
        total_elapsed = time.perf_counter() - overall_start
        error_type = type(e).__name__
        
        print(f"\n❌ TEST FAILED: {error_type}")
        print(f"{'='*70}")
        print(f"Error: {e}")
        import traceback
        print("\n📋 Full traceback:")
        traceback.print_exc()
        print(f"{'='*70}\n")
        
        return jsonify({
            'success': False,
            'error': str(e),
            'error_type': error_type,
            'elapsed_ms': total_elapsed * 1000,
            'diagnostics': {
                'clip_enabled': CLIP_ENABLED,
                'clip_available': CLIP_AVAILABLE
            }
        }), 500


# Admin routes
@app.route('/admin')
@login_required
def admin_dashboard():
    """Admin dashboard - only accessible to authenticated admins"""
    if not is_admin():
        flash('Access denied. Administrator privileges required.', 'error')
        return redirect(url_for('admin_login'))
    return render_template('admin/dashboard.html')

@app.route('/admin/users', methods=['GET', 'POST'])
@login_required
def admin_users():
    """Admin users management page - only accessible to authenticated admins"""
    if not is_admin():
        flash('Access denied. Administrator privileges required.', 'error')
        return redirect(url_for('admin_login'))
    
    # Handle POST requests for user actions
    if request.method == 'POST':
        action = request.form.get('action')
        user_id = request.form.get('user_id')
        
        if action == 'delete_user' and user_id:
            # Delete user
            if delete_user(user_id):
                flash('✅ User deleted successfully.', 'success')
            else:
                flash('❌ Error: Failed to delete user. Please try again.', 'error')
        elif action == 'update_status' and user_id:
            # Update user status
            status = request.form.get('status')
            if status in ['active', 'blocked']:
                if update_user_status(user_id, status):
                    status_text = 'activated' if status == 'active' else 'blocked'
                    flash(f'✅ User status updated: {status_text}.', 'success')
                else:
                    flash('❌ Error: Failed to update user status.', 'error')
            else:
                flash('❌ Invalid status value.', 'error')
        
        # Preserve search and pagination parameters
        search = request.args.get('search', '')
        page = request.args.get('page', 1, type=int)
        return redirect(url_for('admin_users', search=search, page=page))
    
    # GET request - display users with pagination and search
    search_query = request.args.get('search', '').strip()
    page = request.args.get('page', 1, type=int)
    per_page = 15  # Users per page
    
    # Get filtered users
    all_users = get_all_users()
    
    # Apply search filter
    if search_query:
        filtered_users = [
            user for user in all_users
            if search_query.lower() in user['email'].lower() or 
               search_query.lower() in user['status'].lower()
        ]
    else:
        filtered_users = all_users
    
    # Calculate pagination
    total_users = len(filtered_users)
    total_pages = max(1, (total_users + per_page - 1) // per_page)
    page = max(1, min(page, total_pages))  # Ensure page is in valid range
    
    start_idx = (page - 1) * per_page
    end_idx = start_idx + per_page
    paginated_users = filtered_users[start_idx:end_idx]
    
    return render_template(
        'admin/users.html',
        users=paginated_users,
        total_users=total_users,
        page=page,
        total_pages=total_pages,
        per_page=per_page,
        search_query=search_query
    )

# Firebase Authentication routes
@app.route('/auth/firebase/', methods=['POST'])
def firebase_auth_route():
    """Handle Firebase authentication"""
    try:
        data = request.get_json()
        id_token = data.get('token')
        
        if not id_token:
            return jsonify({'error': 'Token is required'}), 400
        
        # Verify Firebase ID token
        try:
            decoded_token = firebase_auth.verify_id_token(id_token)
        except Exception as e:
            print(f"Firebase token verification error: {e}")
            return jsonify({'error': 'Invalid token'}), 401
        
        firebase_uid = decoded_token['uid']
        email = decoded_token.get('email')
        name = decoded_token.get('name', '')
        picture = decoded_token.get('picture', '')
        
        if not email:
            return jsonify({'error': 'Email not found in token'}), 400
        
        # Get or create user
        user_data, error = get_or_create_firebase_user(firebase_uid, email, name, picture)
        if error:
            return jsonify({'error': error}), 500
        
        if user_data['status'] == 'blocked':
            return jsonify({'error': 'Your account has been blocked. Please contact support.'}), 403
        
        # Login user
        user = User(
            user_data['id'],
            email=user_data['email'],
            is_admin_user=False,
            name=user_data.get('name'),
            picture=user_data.get('picture')
        )
        login_user(user, remember=True)
        
        return jsonify({
            'success': True,
            'message': 'Logged in successfully',
            'redirect': url_for('user_dashboard')
        })
        
    except Exception as e:
        print(f"Firebase auth error: {e}")
        return jsonify({'error': 'Authentication failed'}), 500

# User authentication routes
@app.route('/signup', methods=['GET', 'POST'])
def signup():
    """User registration page"""
    if current_user.is_authenticated:
        # If already logged in, redirect to dashboard
        if hasattr(current_user, 'is_admin_user') and current_user.is_admin_user:
            return redirect(url_for('admin_dashboard'))
        return redirect(url_for('user_dashboard'))
    
    if request.method == 'POST':
        email = request.form.get('email', '').strip().lower()
        password = request.form.get('password', '')
        password_confirm = request.form.get('password_confirm', '')
        
        # Validate email
        is_valid_email, email_error = validate_email(email)
        if not is_valid_email:
            flash(f'❌ {email_error}', 'error')
            return render_template('signup.html')
        
        # Validate password
        is_valid_password, password_error = validate_password(password)
        if not is_valid_password:
            flash(f'❌ {password_error}', 'error')
            return render_template('signup.html')
        
        # Check password confirmation
        if password != password_confirm:
            flash('❌ Passwords do not match', 'error')
            return render_template('signup.html')
        
        # Create user
        user_id, error = create_user(email, password)
        if error:
            flash(f'❌ {error}', 'error')
            return render_template('signup.html')
        
        # Auto-login after registration
        user_data = get_user_by_id(user_id)
        if user_data:
            user = User(
                user_data['id'], 
                email=user_data['email'], 
                is_admin_user=False,
                name=user_data.get('name'),
                picture=user_data.get('picture')
            )
            login_user(user)
            flash('✅ Account created successfully! Welcome to AI SlideRush!', 'success')
            return redirect(url_for('user_dashboard'))
    
    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    """User login page"""
    if current_user.is_authenticated:
        # If already logged in, redirect to appropriate dashboard
        if hasattr(current_user, 'is_admin_user') and current_user.is_admin_user:
            return redirect(url_for('admin_dashboard'))
        return redirect(url_for('user_dashboard'))
    
    if request.method == 'POST':
        email = request.form.get('email', '').strip().lower()
        password = request.form.get('password', '')
        
        if not email or not password:
            flash('❌ Please enter both email and password', 'error')
            return render_template('login.html')
        
        # Authenticate user
        user_data, error = authenticate_user(email, password)
        if error:
            flash(f'❌ {error}', 'error')
            return render_template('login.html')
        
        # Login user
        user = User(
            user_data['id'], 
            email=user_data['email'], 
            is_admin_user=False,
            name=user_data.get('name'),
            picture=user_data.get('picture')
        )
        login_user(user, remember=True)
        flash('✅ Logged in successfully!', 'success')
        
        # Redirect to next page or dashboard
        next_page = request.args.get('next')
        if next_page:
            return redirect(next_page)
        return redirect(url_for('user_dashboard'))
    
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    """User logout"""
    logout_user()
    flash('✅ You have been logged out successfully.', 'success')
    return redirect(url_for('index'))

@app.route('/dashboard')
@login_required
def user_dashboard():
    """User dashboard - personal cabinet"""
    # Redirect admins to admin dashboard
    if hasattr(current_user, 'is_admin_user') and current_user.is_admin_user:
        return redirect(url_for('admin_dashboard'))
    
    # Verify user status (prevent blocked users from accessing dashboard)
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute('SELECT status FROM users WHERE id = ?', (current_user.id,))
        user_row = cursor.fetchone()
        conn.close()
        
        if user_row and user_row['status'] == 'blocked':
            logout_user()
            flash('❌ Your account has been blocked. Please contact support.', 'error')
            return redirect(url_for('index'))
    except Exception as e:
        print(f"Error checking user status: {e}")
    
    # Get search, filter and pagination parameters
    search_query = request.args.get('search', '').strip()
    filter_type = request.args.get('type', '').strip()  # Filter by presentation type
    page = request.args.get('page', 1, type=int)
    per_page = 15
    
    # Get user's presentations
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        # Get total count for stats
        cursor.execute(
            'SELECT COUNT(*) as count FROM presentations WHERE user_id = ?',
            (current_user.id,)
        )
        total_presentations = cursor.fetchone()['count']
        
        # Build query with search and type filter
        query = 'SELECT * FROM presentations WHERE user_id = ?'
        params = [current_user.id]
        
        if search_query:
            query += ' AND topic LIKE ?'
            params.append(f'%{search_query}%')
        
        if filter_type and filter_type in PRESENTATION_TYPES:
            query += ' AND presentation_type = ?'
            params.append(filter_type)
        
        query += ' ORDER BY creation_date DESC'
        
        cursor.execute(query, params)
        all_presentations = [dict(row) for row in cursor.fetchall()]
        
        # Add presentation type info to each presentation
        for pres in all_presentations:
            if not pres.get('presentation_type'):
                pres['presentation_type'] = 'business'  # Default for old presentations
            pres['type_info'] = get_presentation_type_info(pres['presentation_type'])
        
        # Get user data (fixed: was calling fetchone() twice, causing user_data to always be None)
        cursor.execute('SELECT * FROM users WHERE id = ?', (current_user.id,))
        user_row = cursor.fetchone()
        user_data = dict(user_row) if user_row else None
        
        conn.close()
    except Exception as e:
        print(f"Error fetching presentations: {e}")
        all_presentations = []
        total_presentations = 0
        user_data = None
    
    # Calculate pagination
    total_filtered = len(all_presentations)
    total_pages = max(1, (total_filtered + per_page - 1) // per_page)
    page = max(1, min(page, total_pages))
    
    start_idx = (page - 1) * per_page
    end_idx = start_idx + per_page
    paginated_presentations = all_presentations[start_idx:end_idx]
    
    return render_template(
        'dashboard.html',
        presentations=paginated_presentations,
        total_presentations=total_presentations,
        user_data=user_data,
        page=page,
        total_pages=total_pages,
        search_query=search_query,
        filter_type=filter_type,
        presentation_types=PRESENTATION_TYPES
    )

@app.route('/presentation/delete', methods=['POST'])
@login_required
def delete_presentation():
    """Delete user's presentation"""
    # Redirect admins
    if hasattr(current_user, 'is_admin_user') and current_user.is_admin_user:
        flash('❌ Admins cannot delete presentations from user dashboard', 'error')
        return redirect(url_for('admin_dashboard'))
    
    presentation_id = request.form.get('presentation_id')
    if not presentation_id:
        flash('❌ Invalid request', 'error')
        return redirect(url_for('user_dashboard'))
    
    try:
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # Verify ownership
        cursor.execute(
            'SELECT * FROM presentations WHERE id = ? AND user_id = ?',
            (presentation_id, current_user.id)
        )
        presentation = cursor.fetchone()
        
        if not presentation:
            conn.close()
            flash('❌ Presentation not found or access denied', 'error')
            return redirect(url_for('user_dashboard'))
        
        # Delete presentation
        cursor.execute('DELETE FROM presentations WHERE id = ?', (presentation_id,))
        conn.commit()
        conn.close()
        
        flash('✅ Presentation deleted successfully', 'success')
    except Exception as e:
        print(f"Error deleting presentation: {e}")
        flash('❌ Error deleting presentation', 'error')
    
    # Preserve search and pagination
    search = request.form.get('search', '')
    page = request.form.get('page', 1, type=int)
    
    return redirect(url_for('user_dashboard', search=search, page=page))

@app.route('/profile/edit', methods=['GET', 'POST'])
@login_required
def edit_profile():
    """Edit user profile"""
    # Redirect admins
    if hasattr(current_user, 'is_admin_user') and current_user.is_admin_user:
        flash('❌ Admins cannot edit profile from user dashboard', 'error')
        return redirect(url_for('admin_dashboard'))
    
    # Get user data
    try:
        conn = sqlite3.connect(DB_PATH)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM users WHERE id = ?', (current_user.id,))
        user_data = dict(cursor.fetchone())
        conn.close()
    except Exception as e:
        print(f"Error fetching user data: {e}")
        user_data = None
    
    if request.method == 'POST':
        # Only allow password change for non-Google users
        if user_data and not user_data.get('google_id'):
            current_password = request.form.get('current_password', '')
            new_password = request.form.get('new_password', '')
            confirm_password = request.form.get('confirm_password', '')
            
            if current_password or new_password or confirm_password:
                # Validate current password
                if not current_password:
                    flash('❌ Please enter your current password', 'error')
                    return render_template('profile_edit.html', user_data=user_data)
                
                # Verify current password
                if not check_password_hash(user_data['password_hash'], current_password):
                    flash('❌ Current password is incorrect', 'error')
                    return render_template('profile_edit.html', user_data=user_data)
                
                # Validate new password
                is_valid, error_msg = validate_password(new_password)
                if not is_valid:
                    flash(f'❌ {error_msg}', 'error')
                    return render_template('profile_edit.html', user_data=user_data)
                
                # Check password confirmation
                if new_password != confirm_password:
                    flash('❌ New passwords do not match', 'error')
                    return render_template('profile_edit.html', user_data=user_data)
                
                # Update password
                try:
                    conn = sqlite3.connect(DB_PATH)
                    cursor = conn.cursor()
                    new_password_hash = generate_password_hash(new_password)
                    cursor.execute(
                        'UPDATE users SET password_hash = ? WHERE id = ?',
                        (new_password_hash, current_user.id)
                    )
                    conn.commit()
                    conn.close()
                    
                    flash('✅ Password updated successfully!', 'success')
                    return redirect(url_for('user_dashboard'))
                except Exception as e:
                    print(f"Error updating password: {e}")
                    flash('❌ Error updating password', 'error')
        else:
            flash('ℹ️ No changes to save', 'error')
    
    return render_template('profile_edit.html', user_data=user_data)

@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login():
    """Admin login page"""
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        # Check if user exists and password is correct
        if username in ADMIN_USERS and check_password_hash(ADMIN_USERS[username]['password_hash'], password):
            user = User(username, is_admin_user=True)
            login_user(user)
            flash('Logged in successfully.', 'success')
            return redirect(url_for('admin_dashboard'))
        else:
            flash('Invalid username or password.', 'error')
    
    return render_template('admin/login.html')

@app.route('/admin/logout')
@login_required
def admin_logout():
    """Admin logout"""
    logout_user()
    flash('You have been logged out.', 'success')
    return redirect(url_for('admin_login'))


# Application entry point - MUST BE AT THE END OF FILE
if __name__ == '__main__':
    print("\n" + "="*60)
    print("  🎨 AI SlideRush - Presentation Service")
    print("="*60)
    
    # Environment Check
    print("\n🔧 CONFIGURATION CHECK:")
    print(f"   OpenAI API Key: {'✅ Configured' if OPENAI_API_KEY else '❌ Missing'}")
    
    # Image Provider Configuration
    print("\n🖼️ IMAGE PROVIDERS:")
    print(f"   Mode: {IMAGE_PROVIDER_MODE.upper()}")
    print(f"   Pexels API: {'✅ Configured' if PEXELS_API_KEY else '⚠️ Missing'}")
    print(f"   Unsplash API: {'✅ Configured' if UNSPLASH_ACCESS_KEY else '⚠️ Not configured (optional)'}")
    
    if IMAGE_PROVIDER_MODE == 'mixed':
        print("   Strategy: Pexels → Unsplash fallback")
    elif IMAGE_PROVIDER_MODE == 'pexels':
        print("   Strategy: Pexels only")
    elif IMAGE_PROVIDER_MODE == 'unsplash':
        print("   Strategy: Unsplash only")
    
    # LibreTranslate / Translation
    print("\n🌍 TRANSLATION:")
    print(f"   Enabled: {TRANSLATION_ENABLED}")
    print(f"   Provider: {TRANSLATION_PROVIDER}")
    if TRANSLATION_ENABLED and TRANSLATION_PROVIDER == 'libre':
        print(f"   LibreTranslate URL: {LIBRETRANSLATE_URL}")
        print(f"   Reachable: {is_libretranslate_available()}")
    elif TRANSLATION_ENABLED and TRANSLATION_PROVIDER == 'external':
        print(f"   External URL: {EXTERNAL_TRANSLATE_URL if EXTERNAL_TRANSLATE_URL else 'Not configured'}")
    
    # Server Start
    port = int(os.environ.get("PORT", 5000))
    print("\n🚀 STARTING SERVER:")
    print(f"   Port: {port}")
    print(f"   URL: http://localhost:{port}")
    print(f"   Debug Mode: True")
    print("\n" + "="*60)
    print("🎉 Server is ready! Press CTRL+C to stop.")
    print("="*60 + "\n")
    
    app.run(debug=True, host='0.0.0.0', port=port)
